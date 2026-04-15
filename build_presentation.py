"""
KYSM Level-5 Roadmap Presentation Generator
Technical audience: Software Engineers / Developers
NRL Brand: "Energy in Motion"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy

# ── Brand Colours ────────────────────────────────────────────────
NRL_YELLOW  = RGBColor(0xFF, 0xD1, 0x00)   # #FFD100
NRL_RED     = RGBColor(0xFF, 0x4A, 0x00)   # #FF4A00
NRL_BLUE    = RGBColor(0x00, 0x3D, 0xA5)   # #003DA5
NRL_DARK    = RGBColor(0x00, 0x2D, 0x72)   # #002D72  (dark navy)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)
GREEN_PASS  = RGBColor(0x00, 0x8A, 0x00)
AMBER_WARN  = RGBColor(0xFF, 0xA5, 0x00)

# ── Helpers ──────────────────────────────────────────────────────

def set_bg(slide, rgb: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, rgb: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    return shape

def accent_bar(slide, left, top, width=9.5, height=0.06, color=NRL_YELLOW):
    """4px yellow accent bar below slide titles"""
    add_rect(slide, left, top, width, height, color)

def add_textbox(slide, text, left, top, width, height,
                font_name="Montserrat", font_size=18,
                bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_textbox(slide, items, left, top, width, height,
                       font_name="Open Sans", font_size=16,
                       color=NEAR_BLACK, bullet_color=NRL_YELLOW):
    """Add a text box with multiple bullet items"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def add_slide_number(slide, number, total):
    """Slide number bottom-right"""
    add_textbox(slide, f"Slide {number} of {total}",
                7.5, 7.1, 2.0, 0.4,
                font_size=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

def add_nrl_tag(slide, text="Numaligarh Refinery Limited | Internal Use Only"):
    """Footer tag bottom-left"""
    add_textbox(slide, text,
                0.3, 7.1, 5.0, 0.4,
                font_size=9, color=DARK_GRAY, align=PP_ALIGN.LEFT)

def add_logo_placeholder(slide):
    """Yellow accent square in top-right as logo placeholder"""
    shape = slide.shapes.add_shape(
        1, Inches(8.9), Inches(0.2), Inches(0.65), Inches(0.65)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NRL_YELLOW
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NRL"
    run.font.name = "Montserrat"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NRL_DARK

# ── Slide Factories ──────────────────────────────────────────────

def cover_slide(prs, slide_num, total):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, NRL_DARK)

    # Large diagonal accent
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(3.5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x1A, 0x45)
    shape.line.fill.background()

    # Yellow accent stripe
    add_rect(slide, 0, 4.8, 13.33, 0.08, NRL_YELLOW)

    # Logo top-left
    logo = slide.shapes.add_shape(1, Inches(0.4), Inches(0.4), Inches(1.2), Inches(1.2))
    logo.fill.solid()
    logo.fill.fore_color.rgb = NRL_YELLOW
    logo.line.fill.background()
    tf = logo.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NRL"
    run.font.name = "Montserrat"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = NRL_DARK

    # Title
    add_textbox(slide, "Know Your SAP Masters",
                0.5, 1.8, 12.0, 1.2,
                font_name="Montserrat", font_size=44,
                bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, "Level-5 Technical Roadmap",
                0.5, 3.0, 12.0, 0.7,
                font_name="Montserrat", font_size=26,
                bold=False, color=NRL_YELLOW, align=PP_ALIGN.LEFT)

    # Tagline
    add_textbox(slide, "5-Pillar RAG • 13-Phase Execution • Multi-Agent Swarm",
                0.5, 3.75, 12.0, 0.5,
                font_name="Open Sans", font_size=16,
                bold=False, color=RGBColor(0xCC, 0xD6, 0xFF), align=PP_ALIGN.LEFT)

    # Meta line
    add_textbox(slide, "April 15, 2026  |  For Technical Audience",
                0.5, 5.2, 8.0, 0.4,
                font_name="Open Sans", font_size=13,
                bold=False, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT)

    add_textbox(slide, "Energy in Motion",
                9.5, 6.8, 3.5, 0.4,
                font_name="Montserrat", font_size=13,
                bold=False, color=NRL_YELLOW, align=PP_ALIGN.RIGHT)

    return slide

def section_divider(prs, title, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, NRL_DARK)

    # Yellow accent bar
    add_rect(slide, 0, 3.3, 13.33, 0.07, NRL_YELLOW)

    add_textbox(slide, title,
                0.5, 2.3, 12.0, 1.2,
                font_name="Montserrat", font_size=40,
                bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_textbox(slide, subtitle,
                    0.5, 3.55, 12.0, 0.6,
                    font_name="Open Sans", font_size=18,
                    bold=False, color=NRL_YELLOW, align=PP_ALIGN.LEFT)

    add_logo_placeholder(slide)
    return slide

def content_slide(prs, slide_num, total, title, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)

    # Header bar
    add_rect(slide, 0, 0, 13.33, 0.85, NRL_DARK)
    # Yellow accent strip below header
    add_rect(slide, 0, 0.85, 13.33, 0.055, NRL_YELLOW)

    # Title
    add_textbox(slide, title,
                0.35, 0.12, 12.5, 0.65,
                font_name="Montserrat", font_size=26,
                bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        add_textbox(slide, subtitle,
                    0.35, 0.93, 12.0, 0.4,
                    font_name="Open Sans", font_size=13,
                    bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT)
        accent_bar(slide, 0.35, 1.32, 12.5, 0.04, NRL_YELLOW)
    else:
        accent_bar(slide, 0.35, 0.93, 12.5, 0.04, NRL_YELLOW)

    add_nrl_tag(slide)
    add_logo_placeholder(slide)
    add_slide_number(slide, slide_num, total)
    return slide

def two_column_slide(prs, slide_num, total, title,
                     left_title, left_items,
                     right_title, right_items,
                     subtitle=""):
    slide = content_slide(prs, slide_num, total, title, subtitle)

    # Left column header
    add_rect(slide, 0.35, 1.45, 5.8, 0.4, NRL_BLUE)
    add_textbox(slide, left_title,
                0.35, 1.45, 5.8, 0.4,
                font_name="Montserrat", font_size=14,
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left content
    y = 1.95
    for item in left_items:
        add_textbox(slide, f"• {item}",
                    0.35, y, 5.8, 0.38,
                    font_name="Open Sans", font_size=13,
                    bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT)
        y += 0.42

    # Right column header
    add_rect(slide, 6.55, 1.45, 6.4, 0.4, NRL_BLUE)
    add_textbox(slide, right_title,
                6.55, 1.45, 6.4, 0.4,
                font_name="Montserrat", font_size=14,
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right content
    y = 1.95
    for item in right_items:
        add_textbox(slide, f"• {item}",
                    6.55, y, 6.4, 0.38,
                    font_name="Open Sans", font_size=13,
                    bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT)
        y += 0.42

    return slide

def table_slide(prs, slide_num, total, title, headers, rows,
                col_widths=None, subtitle=""):
    slide = content_slide(prs, slide_num, total, title, subtitle)

    if col_widths is None:
        col_widths = [1.5] * len(headers)

    # Header row
    x = 0.35
    y = 1.45
    row_h = 0.38
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, y, w, row_h, NRL_BLUE)
        add_textbox(slide, h, x + 0.05, y, w - 0.1, row_h,
                    font_name="Montserrat", font_size=11,
                    bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        x += w

    # Data rows
    y += row_h
    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        x = 0.35
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, w, row_h, bg)
            # Status colouring
            cell_color = NEAR_BLACK
            if cell == "✅ Working" or cell == "✅ COMPLETE":
                cell_color = GREEN_PASS
            elif cell == "🚧 Pending":
                cell_color = AMBER_WARN
            add_textbox(slide, str(cell), x + 0.05, y, w - 0.1, row_h,
                        font_name="Open Sans", font_size=11,
                        bold=False, color=cell_color, align=PP_ALIGN.LEFT)
            x += w
        y += row_h

    return slide

def status_grid_slide(prs, slide_num, total, title, items, subtitle=""):
    """4-column status grid slide"""
    slide = content_slide(prs, slide_num, total, title, subtitle)

    cols = 4
    box_w = 2.95
    box_h = 1.4
    start_x = 0.35
    gap_x = 0.12
    start_y = 1.55

    for idx, (name, status, detail) in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (box_w + gap_x)
        y = start_y + row * (box_h + 0.18)

        # Box background
        box_color = NRL_BLUE if "✅" in status else (NRL_RED if "❌" in status else AMBER_WARN)
        add_rect(slide, x, y, box_w, box_h, box_color)

        # Component name
        add_textbox(slide, name, x + 0.1, y + 0.08, box_w - 0.2, 0.35,
                    font_name="Montserrat", font_size=13,
                    bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # Status
        add_textbox(slide, status, x + 0.1, y + 0.45, box_w - 0.2, 0.28,
                    font_name="Open Sans", font_size=12,
                    bold=False, color=NRL_YELLOW, align=PP_ALIGN.LEFT)

        # Detail
        add_textbox(slide, detail, x + 0.1, y + 0.75, box_w - 0.2, 0.55,
                    font_name="Open Sans", font_size=10,
                    bold=False, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.LEFT)

    return slide

def bullet_slide(prs, slide_num, total, title, items, subtitle="", two_col=False,
                 left_items=None, right_items=None):
    if two_col and left_items and right_items:
        return two_column_slide(prs, slide_num, total, title,
                                title.split("|")[0].strip() if "|" in title else "Architecture",
                                left_items,
                                title.split("|")[-1].strip() if "|" in title else "Components",
                                right_items, subtitle)
    else:
        slide = content_slide(prs, slide_num, total, title, subtitle)
        y = 1.55
        for item in items:
            add_textbox(slide, f"• {item}",
                        0.35, y, 12.5, 0.4,
                        font_name="Open Sans", font_size=14,
                        bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT)
            y += 0.46
        return slide

# ── Build Presentation ───────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 14  # total slides (for page numbers)

# ── Slide 1: Cover ────────────────────────────────────────────────
cover_slide(prs, 1, TOTAL)

# ── Slide 2: Agenda ───────────────────────────────────────────────
agenda_items = [
    "5-Pillar RAG Architecture — what it is and how it works",
    "13-Phase Execution Roadmap — every phase from meta-path to result masking",
    "Infrastructure Stack — Qdrant, Memgraph, Redis, RabbitMQ, Celery",
    "Memgraph Migration M1–M7 — from in-memory NetworkX to distributed Cypher",
    "Harness Engineering — Deep Harnessing, Memory Compounding, Threat Sentinel",
    "Multi-Agent Domain Swarm — 7 domain agents, Planner, Synthesis, Message Bus",
    "Pending Work — Real SAP HANA, BAPI Workflows, Benchmark suite",
    "Key Files & Code References — for engineers joining the project",
]
section_divider(prs, "Agenda", "What we will cover in this presentation")
slide = content_slide(prs, 2, TOTAL, "Agenda", "8 sections — ~15 minutes")
y = 1.55
for item in agenda_items:
    add_textbox(slide, f"  ▸  {item}",
                0.35, y, 12.5, 0.42,
                font_name="Open Sans", font_size=14,
                bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT)
    y += 0.48

# ── Slide 3: 5-Pillar RAG Architecture ──────────────────────────
headers = ["Pillar", "Name", "Component", "Status"]
rows = [
    ["1", "Role-Aware Security", "security.py — SAPAuthContext, row/column masking", "✅ Working"],
    ["2", "Agentic Orchestrator", "orchestrator.py — 10-step run_agent_loop()", "✅ Working"],
    ["3", "Schema RAG", "Qdrant + ChromaDB, schema_lookup()", "✅ Working"],
    ["4", "SQL Pattern RAG", "sql_pattern_lookup(), 68+ patterns, 18 domains", "✅ Working"],
    ["5", "Graph RAG", "NetworkX FK graph, AllPathsExplorer, TemporalGraphRAG, Meta-Path Library", "✅ Working"],
    ["5½", "Graph Embedding Search", "Node2Vec structural + text hybrid, ChromaDB collections", "✅ Working"],
]
table_slide(prs, 3, TOTAL, "5-Pillar RAG Architecture",
            headers, rows,
            col_widths=[0.6, 1.8, 7.8, 1.1],
            subtitle="Core security + orchestration + retrieval pipeline")

# ── Slide 4: 13-Phase Execution Roadmap ─────────────────────────
headers = ["Phase", "Name", "Description", "Status"]
rows = [
    ["0", "Meta-Path Match", "Fast-path template matching (14 pre-computed JOIN paths)", "✅ Working"],
    ["1", "Schema RAG", "Qdrant semantic search over DDIC metadata", "✅ Working"],
    ["1.5", "Graph Embedding", "Node2Vec structural + text hybrid table discovery", "✅ Working"],
    ["1.75", "QM Semantic Search", "20yr QM notification long-text semantic search", "✅ Working"],
    ["2", "SQL Pattern RAG", "Qdrant proven SQL patterns (18 domains)", "✅ Working"],
    ["2b", "Temporal Detection", "Date/fiscal anchor detection → temporal filters", "✅ Working"],
    ["5.5", "Validation Harness", "SELECT COUNT(*) dry-run → syntax validation → autonomous fix", "✅ NEW"],
    ["6", "Self-Healing", "Rule-based SQL correction (10 error codes → 6 heal strategies)", "✅ Working"],
    ["6b", "Memory Compounding", "Auto-vectorize healed SQL back into Qdrant pattern store", "✅ NEW"],
    ["6c", "Threat Sentinel", "6 threat engines + dynamic AuthContext tightening", "✅ NEW"],
    ["10", "Multi-Agent Swarm", "Planner + 7 Domain Agents + Synthesis Agent", "✅ LIVE"],
    ["13", "Inter-Agent Message Bus", "Redis pub/sub + streams, 6 message types", "✅ IMPLEMENTED"],
]
table_slide(prs, 4, TOTAL, "13-Phase Execution Roadmap (condensed)",
            headers, rows,
            col_widths=[0.65, 1.5, 8.2, 1.95],
            subtitle="Full pipeline: query → meta-path → schema → SQL → graph → assembly → validation → execution → masking")

# ── Slide 5: Infrastructure Status ───────────────────────────────
infra_items = [
    ("Qdrant", "✅ ACTIVE", "4 collections: sap_schema, sql_patterns, graph_node_embeddings, graph_table_context"),
    ("Memgraph", "✅ ACTIVE", "114 nodes / 137 edges. All NetworkX edges synced. Cypher queries working."),
    ("ChromaDB", "🔄 RETIRED", "All collections migrated to Qdrant. Schema RAG, Pattern RAG, QM, Graph Embeddings."),
    ("RabbitMQ", "✅ ACTIVE", "amqp://sapmasters:sapmasters123@localhost:5672//"),
    ("Redis", "✅ ACTIVE", "localhost:6379/0 — dialog sessions, role preferences, memory"),
    ("Celery Worker", "✅ ACTIVE", "4 threads, queues: agent + priority. Swarm autoscaling via domain_tasks.py"),
]
status_grid_slide(prs, 5, TOTAL, "Infrastructure Status", infra_items,
                  subtitle="As of April 15, 2026 — 18:00 IST")

# ── Slide 6: Memgraph Migration ──────────────────────────────────
headers = ["Phase", "Description", "Status"]
rows = [
    ["M1", "Memgraph 2.12.0 + Lab — Docker Compose", "✅ Complete"],
    ["M1", "Schema init + load (init_schema.cql) — 114 nodes, 47 edges", "✅ Complete"],
    ["M2", "Cypher port — native pattern matching (no LENGTH/relationships on paths)", "✅ Complete"],
    ["M3", "use_memgraph flag in main.py + auto-sync on startup", "✅ Complete"],
    ["M4", "Qdrant cluster migration (Schema RAG + Pattern RAG + QM + Graph Embeddings)", "✅ Complete"],
    ["M5", "Celery async worker pool (RabbitMQ + Redis + 4-thread worker)", "✅ Complete"],
    ["M6", "Load testing + production tuning (memgraph_load_test.py)", "🚧 Pending"],
    ["M7", "Real SAP HANA connection (hana_pool.py, HANA_MODE=pool)", "🚧 Pending"],
]
table_slide(prs, 6, TOTAL, "Memgraph Migration (M1–M7)",
            headers, rows,
            col_widths=[0.65, 9.5, 2.15],
            subtitle="In-memory NetworkX → distributed Memgraph Cypher. Non-breaking migration.")

# ── Slide 7: Harness Engineering ────────────────────────────────
two_column_slide(prs, 7, TOTAL,
    "Harness Engineering — Three Innovations",
    "Phase 5.5: Deep Harnessing",
    [
        "File: sql_executor.py + orchestrator.py",
        "SELECT COUNT(*) dry-run before real execution",
        "Error codes: 37000, ORA-01476, ORA-00942, ORA-01799",
        "On failure: SelfHealer.heal() → re-test → then execute",
        "Zero human intervention for syntax recovery",
    ],
    "Phase 6b: Memory Compounding",
    [
        "File: orchestrator.py — Step 8b",
        "Self-heal success → build intent string",
        "store_manager.load_domain() → Qdrant upsert",
        "AI autonomously expands pattern library",
        "Healed queries become future fast-path hits",
    ],
    subtitle="Self-validated, self-improving agent loop")

two_column_slide(prs, 7, TOTAL,
    "Harness Engineering — Phase 6c: Threat Sentinel",
    "6 Threat Engines (security_sentinel.py)",
    [
        "CROSS_MODULE_ESCALATION — multi-hop graph traversal to out-of-scope tables",
        "SCHEMA_ENUMERATION — bulk table discovery probes (>5 new tables/query)",
        "DENIED_TABLE_PROBE — repeated attempts to access blocked tables",
        "DATA_EXFILTRATION — unusually large result sets (>5,000 rows)",
        "TEMPORAL_INFERENCE — restricted historical period access by HR/AP roles",
        "ROLE_IMPERSONATION — sudden cross-domain shift mid-session",
    ],
    "Sentinel Actions",
    [
        "DISABLED | AUDIT | ENFORCING modes",
        "Dynamic AuthContext tightening (tightness 0→3)",
        "Denied table expansion, masked field expansion",
        "SIEM/webhook alerts on ENFORCING violations",
        "Pre-execution gate in orchestrator",
        "Verdict returned in API 'sentinel' key",
    ],
    subtitle="Proactive threat detection — 32KB security_sentinel.py")

# Actually redo slide 7 as a single slide
# Let me make a proper slide 7
slide7 = content_slide(prs, 7, TOTAL, "Harness Engineering — Three Innovations",
                       "Self-validated, self-improving agent loop")
add_textbox(slide7, "Phase 5.5 — Deep Harnessing (Validation Harness)",
            0.35, 1.45, 6.0, 0.35,
            font_name="Montserrat", font_size=14, bold=True, color=NRL_BLUE)
items_55 = [
    "SELECT COUNT(*) dry-run before real execution",
    "Error codes: 37000, ORA-01476, ORA-00942, ORA-01799",
    "On failure: SelfHealer.heal() → re-test → execute",
    "Zero human intervention for syntax recovery",
]
y = 1.85
for it in items_55:
    add_textbox(slide7, f"  • {it}", 0.35, y, 6.0, 0.33,
                font_name="Open Sans", font_size=12, color=NEAR_BLACK)
    y += 0.37

add_textbox(slide7, "Phase 6b — Memory Compounding",
            6.7, 1.45, 6.2, 0.35,
            font_name="Montserrat", font_size=14, bold=True, color=NRL_BLUE)
items_6b = [
    "Self-heal success → Qdrant upsert",
    "AI autonomously expands pattern library",
    "Healed queries become future fast-path hits",
]
y = 1.85
for it in items_6b:
    add_textbox(slide7, f"  • {it}", 6.7, y, 6.2, 0.33,
                font_name="Open Sans", font_size=12, color=NEAR_BLACK)
    y += 0.37

# Phase 6c
add_rect(slide7, 0.35, 3.45, 12.5, 0.38, NRL_DARK)
add_textbox(slide7, "Phase 6c — Proactive Threat Sentinel (32KB security_sentinel.py)",
            0.4, 3.45, 12.3, 0.38,
            font_name="Montserrat", font_size=13, bold=True, color=WHITE)
threat_items = [
    "6 engines: CROSS_MODULE_ESCALATION, SCHEMA_ENUMERATION, DENIED_TABLE_PROBE, DATA_EXFILTRATION, TEMPORAL_INFERENCE, ROLE_IMPERSONATION",
    "Modes: DISABLED | AUDIT | ENFORCING",
    "Actions: AuthContext tightening → denied_tables expansion → SIEM/webhook alerts",
    "Pre-execution gate in orchestrator; verdict returned in API 'sentinel' key",
]
y = 3.93
for it in threat_items:
    add_textbox(slide7, f"  • {it}", 0.35, y, 12.5, 0.34,
                font_name="Open Sans", font_size=12, color=NEAR_BLACK)
    y += 0.38

# ── Slide 8: Multi-Agent Domain Swarm ────────────────────────────
headers = ["Component", "File", "Status"]
rows = [
    ["7 Domain Agents (pur/bp/mm/qm/sd/wm/cross)", "domain_agents.py", "✅ Working"],
    ["Planner Agent + Complexity Analyzer", "swarm/planner_agent.py (19KB)", "✅ LIVE"],
    ["Synthesis Agent (merge + rank + conflicts)", "swarm/synthesis_agent.py (16KB)", "✅ LIVE"],
    ["Swarm entry point", "swarm/__init__.py (2KB)", "✅ LIVE"],
    ["Orchestrator use_swarm flag + API wiring", "orchestrator.py, api/endpoints/chat.py", "✅ LIVE"],
    ["Frontend default use_swarm=True + swarm UI", "frontend/app.py", "✅ LIVE"],
    ["Inter-Agent Message Bus (Redis pub/sub)", "app/core/message_bus.py", "✅ IMPLEMENTED"],
    ["Negotiation Protocol (4-phase)", "app/core/negotiation_protocol.py", "✅ IMPLEMENTED"],
    ["Swarm Autoscaling — Per-domain Celery queues", "app/workers/domain_tasks.py", "✅ IMPLEMENTED"],
]
table_slide(prs, 8, TOTAL, "Multi-Agent Domain Swarm — Architecture",
            headers, rows,
            col_widths=[3.5, 5.5, 2.3],
            subtitle="Port 8001 live — 7 specialist agents + Planner + Synthesis + Message Bus + Negotiation")

# ── Slide 9: Swarm Bug Fixes ─────────────────────────────────────
slide9 = content_slide(prs, 9, TOTAL, "Swarm Bugs Fixed During Activation",
                       "3 bugs resolved before go-live")
bugs = [
    ("Bug 1", "tables_involved referenced before init",
     "Added early init before sentinel evaluation gate in orchestrator.py"),
    ("Bug 2", "cross_agent list index out of range",
     "_mask_results guard: table = self.primary_tables[0] if self.primary_tables else ''"),
    ("Bug 3", "abs(min(vals), 0.01) Python syntax error",
     "Fixed to max(abs(min(vals)), 0.01) in synthesis_agent.py"),
    ("Bug 4", "UnboundLocalError on get_harness_runs",
     "Removed inner 'from ... import' inside try blocks in orchestrator.py"),
    ("Bug 5", "IndexError: list index out of range (tables_involved)",
     "Added elif tables_involved / else: base_sql='' guard at orchestrator.py ~1453"),
    ("Bug 6", "IndexError: decision.assignments[0] empty",
     "Added if not decision.assignments guard in planner_agent.py dispatch_single"),
]
y = 1.5
for bug_id, bug_title, fix in bugs:
    add_rect(slide9, 0.35, y, 0.9, 0.4, NRL_BLUE)
    add_textbox(slide9, bug_id, 0.35, y, 0.9, 0.4,
                font_name="Montserrat", font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide9, bug_title, 1.35, y, 3.2, 0.4,
                font_name="Montserrat", font_size=11, bold=True, color=NRL_RED)
    add_textbox(slide9, fix, 4.65, y, 8.2, 0.4,
                font_name="Open Sans", font_size=11, bold=False, color=NEAR_BLACK)
    y += 0.56

# ── Slide 10: Pending Work ───────────────────────────────────────
p_items = [
    ("P0 — Real SAP HANA", "Wire hdbcli to replace mock executor. Final production barrier."),
    ("P1 — BAPI Workflow Harness", "BAPI_PO_CHANGE, BAPI_VENDOR_CREATE, BAPI_MATERIAL_SAVEDATA, BAPI_SALESORDER_CHANGE"),
    ("P2 — Inter-Agent Message Bus", "Break domain agents out of star-topology via direct agent-to-agent negotiation"),
    ("P3 — 50-Query Benchmark", "Golden dataset to feed Eval Alerting with real failure signals"),
    ("M6 — Load Testing", "Run memgraph_load_test.py at concurrency 10–20, tune pool_size=20+"),
    ("M7 — Memgraph M2", "Replace NetworkX with direct Memgraph Cypher queries (find_all_ranked_paths_native)"),
]
slide10 = content_slide(prs, 10, TOTAL, "Pending Work — P0 to P3 + M6–M7",
                        "Items blocking production readiness")
y = 1.55
for p_id, desc in p_items:
    add_rect(slide10, 0.35, y, 2.1, 0.45, NRL_DARK)
    add_textbox(slide10, p_id, 0.35, y, 2.1, 0.45,
                font_name="Montserrat", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide10, desc, 2.55, y, 10.4, 0.45,
                font_name="Open Sans", font_size=13, bold=False, color=NEAR_BLACK)
    y += 0.58

# ── Slide 11: Key Files Reference ───────────────────────────────
headers = ["File", "Purpose", "Phase"]
rows = [
    ["orchestrator.py", "Main agentic loop — 10-step run_agent_loop()", "Pillar 2"],
    ["security.py", "SAPAuthContext + SecurityMesh", "Pillar 1"],
    ["security_sentinel.py", "Proactive Threat Sentinel (32KB)", "Phase 6c"],
    ["self_healer.py", "SQL self-healing engine", "Phase 6"],
    ["vector_store.py", "Dual-backend Qdrant + ChromaDB", "Pillar 3"],
    ["graph_embedding_store.py", "Node2Vec + text hybrid", "Pillar 5½"],
    ["meta_path_library.py", "14 meta-paths, 22+ JOIN variants", "Phase 0"],
    ["graph_store.py", "NetworkX FK graph + AllPathsExplorer + TemporalGraphRAG", "Pillar 5"],
    ["sql_executor.py", "SAP HANA executor (mock + hdbcli)", "Phase 7"],
    ["orchestrator_tools.py", "Tool registry + 12 tool implementations", "Orchestration"],
    ["planner_agent.py", "Planner Agent + Complexity Analyzer (19KB)", "Phase 10"],
    ["synthesis_agent.py", "Synthesis Agent (16KB)", "Phase 10"],
    ["domain_tasks.py", "Per-domain Celery tasks + queue routing", "Phase SWARM-AUTO"],
    ["message_bus.py", "Redis pub/sub + streams", "Phase 13"],
    ["negotiation_protocol.py", "4-phase conflict resolution", "Phase 13"],
]
table_slide(prs, 11, TOTAL, "Key Files Reference",
            headers, rows,
            col_widths=[3.0, 7.5, 1.8],
            subtitle="Core engineering reference — all files under backend/app/")

# ── Slide 12: Swarm Autoscaling Detail ──────────────────────────
slide12 = content_slide(prs, 12, TOTAL, "Swarm Autoscaling — Per-Domain Celery Queues",
                        "Horizontal scaling via independent domain workers")
items = [
    "7 domain-specific queues: pur_queue, bp_queue, mm_queue, qm_queue, sd_queue, wm_queue, cross_queue",
    "dispatch_domain_group(domain, tasks) — dispatches domain tasks to Celery",
    "collect_group_results(async_results) — threads results back to Synthesis Agent",
    "domain_tasks.py (381 lines) — replaces ThreadPoolExecutor in dispatch_parallel",
    "Celery task per domain agent = independent horizontal scaling per domain",
    "agent queue remains as fallback for non-domain tasks",
    "Broker: RabbitMQ (amqp://sapmasters:sapmasters123@localhost:5672//)",
    "Result backend: Redis (localhost:6379/0)",
    "Worker: 4 threads, queues: agent + priority + all domain queues",
]
y = 1.55
for item in items:
    add_textbox(slide12, f"  ◆ {item}",
                0.35, y, 12.5, 0.4,
                font_name="Open Sans", font_size=13, bold=False, color=NEAR_BLACK)
    y += 0.47

# ── Slide 13: API Contract & Usage ──────────────────────────────
slide13 = content_slide(prs, 13, TOTAL, "API Contract & Usage",
                        "How engineers should call the orchestrator")
add_textbox(slide13, "POST /api/v1/chat/master-data",
            0.35, 1.5, 12.5, 0.38,
            font_name="Montserrat", font_size=15, bold=True, color=NRL_BLUE)
add_textbox(slide13, "Request Body (JSON)",
            0.35, 1.93, 4.0, 0.35,
            font_name="Montserrat", font_size=12, bold=True, color=NEAR_BLACK)
code_body = '{\n  "query": "vendor payment terms for company code 1000",\n  "user_role": "AP_CLERK",\n  "domain": "auto",\n  "use_swarm": False\n}'
add_rect(slide13, 0.35, 2.3, 5.8, 1.3, RGBColor(0xF0, 0xF0, 0xF0))
add_textbox(slide13, code_body, 0.45, 2.32, 5.6, 1.26,
            font_name="Courier New", font_size=11, bold=False, color=NEAR_BLACK)

add_textbox(slide13, "Response Fields (key)",
            6.55, 1.93, 6.2, 0.35,
            font_name="Montserrat", font_size=12, bold=True, color=NEAR_BLACK)
resp_fields = [
    "tables_used — list of SAP table names",
    "executed_sql — generated + validated SQL",
    "data — result records (masked per role)",
    "confidence_score — 6-signal composite",
    "tool_trace — phase-by-phase execution log",
    "sentinel — threat detection verdict",
    "swarm_routing — single/parallel/cross_module",
]
y = 2.32
for f in resp_fields:
    add_textbox(slide13, f"  • {f}", 6.55, y, 6.2, 0.28,
                font_name="Open Sans", font_size=11, color=NEAR_BLACK)
    y += 0.3

add_textbox(slide13, "use_swarm=True triggers Multi-Agent Domain Swarm (7 agents in parallel)",
            0.35, 4.0, 12.5, 0.35,
            font_name="Open Sans", font_size=12, bold=False, color=NRL_RED)

# ── Slide 14: Thank You ──────────────────────────────────────────
blank_layout = prs.slide_layouts[6]
slide14 = prs.slides.add_slide(blank_layout)
set_bg(slide14, NRL_DARK)

# Diagonal
add_rect(slide14, 0, 0, 4.0, 7.5, RGBColor(0x00, 0x1A, 0x45))

# Yellow accent
add_rect(slide14, 0, 4.5, 13.33, 0.08, NRL_YELLOW)

# Logo
logo = slide14.shapes.add_shape(1, Inches(0.4), Inches(0.4), Inches(1.2), Inches(1.2))
logo.fill.solid()
logo.fill.fore_color.rgb = NRL_YELLOW
logo.line.fill.background()
tf = logo.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "NRL"
run.font.name = "Montserrat"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = NRL_DARK

add_textbox(slide14, "Thank You",
            0.5, 2.0, 12.0, 1.2,
            font_name="Montserrat", font_size=48,
            bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide14, "Questions? Let's discuss.",
            0.5, 3.25, 12.0, 0.6,
            font_name="Open Sans", font_size=22,
            bold=False, color=NRL_YELLOW, align=PP_ALIGN.LEFT)
add_textbox(slide14, "Know Your SAP Masters — Level-5 Technical Roadmap",
            0.5, 3.9, 12.0, 0.4,
            font_name="Open Sans", font_size=14,
            bold=False, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT)
add_textbox(slide14, "Energy in Motion",
            9.5, 6.8, 3.5, 0.4,
            font_name="Montserrat", font_size=13,
            bold=False, color=NRL_YELLOW, align=PP_ALIGN.RIGHT)
add_textbox(slide14, "Numaligarh Refinery Limited  |  Internal Use Only",
            0.5, 6.8, 7.0, 0.4,
            font_size=11, color=RGBColor(0x88, 0x99, 0xAA))

# ── Save ────────────────────────────────────────────────────────
out_path = "C:/Users/vishnu/.openclaw/workspace/SAP_HANA_LLM_VendorChatbot/docs/KYSM_Level5_Roadmap_Technical.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")