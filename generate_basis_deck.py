"""
Know Your SAP Masters — Deck 1: Basis Team (Technical Deep-Dive)
SAP Joule vs. Know Your SAP Masters — Feature Alignment
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Palette ──────────────────────────────────────────────
C_DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x3E)
C_NAVY        = RGBColor(0x1A, 0x37, 0x6E)
C_ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xC2)
C_ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF5)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF7)
C_MID_GRAY    = RGBColor(0x8A, 0x9B, 0xA8)
C_TEXT_DARK   = RGBColor(0x1C, 0x2D, 0x4A)
C_GREEN       = RGBColor(0x00, 0xB0, 0x6A)
C_PURPLE      = RGBColor(0x6B, 0x47, 0xD0)
C_ORANGE      = RGBColor(0xE8, 0x7D, 0x1A)
C_RED         = RGBColor(0xCC, 0x29, 0x2A)
C_TEAL        = RGBColor(0x00, 0x96, 0x88)

PILLAR_COLORS = [C_GREEN, C_PURPLE, C_ORANGE, C_RED, C_TEAL]
PILLAR_NAMES  = ["P1 Role-Aware", "P2 Agentic", "P3 Schema", "P4 SQL", "P5 Graph"]

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, lw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = Pt(lw) if lw else 0
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, sz=11, bold=False, italic=False,
        color=C_TEXT_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.1), fill_rgb=C_DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), W, Pt(4), fill_rgb=C_ACCENT_BLUE)
    txt(slide, title, 0.4, 0.12, W - 0.8, 0.65, sz=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.72, W - 0.8, 0.35, sz=13, color=C_ACCENT_BLUE)


def footer(slide, text="Know Your SAP Masters — Basis Team Deck  |  Confidential"):
    add_rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill_rgb=C_DARK_NAVY)
    txt(slide, text, 0.3, H - Inches(0.28), W - 0.6, 0.26, sz=8, color=C_MID_GRAY)


def section(slide, num, title, subtitle=""):
    add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
    txt(slide, f"0{num}", 0.5, 1.2, 4, 3.5, sz=160, bold=True,
        color=RGBColor(0x1A, 0x37, 0x6E))
    add_rect(slide, 4.5, 2.9, 0.06, 1.8, fill_rgb=C_ACCENT_BLUE)
    txt(slide, title, 4.8, 2.7, 8, 1.2, sz=38, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 4.8, 4.0, 8, 0.6, sz=16, color=C_ACCENT_BLUE)


def table_style(tbl, hdr_rgb, row_rgb=None):
    for ri, row in enumerate(tbl.rows):
        for cell in row.cells:
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = hdr_rgb
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.bold = True; run.font.color.rgb = C_WHITE; run.font.size = Pt(10)
            else:
                cell.fill.fore_color.rgb = row_rgb or C_LIGHT_GRAY
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(10); run.font.color.rgb = C_TEXT_DARK


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
add_rect(slide, 9.2, 0, 4.13, H, fill_rgb=C_NAVY)
add_rect(slide, 9.13, 0, Pt(6), H, fill_rgb=C_ACCENT_BLUE)

lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 1.5, 2.8, 0.42)
lbl.fill.solid(); lbl.fill.fore_color.rgb = C_ACCENT_BLUE; lbl.line.fill.background()
tf = lbl.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "BASIS TEAM — TECHNICAL DECK"
r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C_WHITE

txt(slide, "Know Your\nSAP Masters", 0.5, 2.1, 8.4, 2.2, sz=52, bold=True, color=C_WHITE)
txt(slide, "5-Pillar RAG Architecture — Technical Deep-Dive",
    0.5, 4.35, 8.4, 0.55, sz=20, color=C_ACCENT_BLUE)
add_rect(slide, 0.5, 5.05, 3.5, Pt(3), fill_rgb=C_ACCENT_GOLD)
txt(slide, "For SAP Basis, Data & Integration Teams",
    0.5, 5.2, 8.4, 0.45, sz=13, italic=True, color=C_MID_GRAY)

for i, (name, col) in enumerate(zip(PILLAR_NAMES, PILLAR_COLORS)):
    y = 1.0 + i * 1.2
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 9.55, y, 3.3, 0.85)
    pill.fill.solid(); pill.fill.fore_color.rgb = col; pill.line.fill.background()
    tf2 = pill.text_frame; tf2.word_wrap = False; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run()
    r2.text = name; r2.font.size = Pt(9); r2.font.bold = True
    r2.font.color.rgb = C_WHITE

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM: WHY GENERIC LLMs FAIL ON SAP
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "The Problem: Generic LLMs Fail on SAP S/4HANA",
       "5 failure modes that require domain-specific architecture")

problems = [
    ("90,000+ DDIC Tables",
     "Schema RAG is mandatory — natural language cannot reliably map 'vendor' to LFA1 or 'pricing' to KONP without vector retrieval."),
    ("SAP SQL Idioms",
     "MANDT filters, LOEVM/LVORM deletion flags, ADD_DAYS(ZFBDT,ZBD3T) for overdue dates — LLMs hallucinate these without few-shot grounding."),
    ("Cross-Module Joins",
     "'Which vendors supply Cost Center X?' requires MARA→EINA→LFA1→CSKS traversal — single-shot RAG misses this."),
    ("Authorization at SQL Level",
     "SAP's object-based security (BUKRS, WERKS, EKORG) must be a hard SQL WHERE constraint — not a prompt instruction."),
    ("No Audit Trail",
     "Every generated SQL must be logged with user identity, SQL hash, auth context, and result rows for SOX compliance."),
]

for i, (title, desc) in enumerate(problems):
    y = 1.35 + i * 1.18
    add_rect(slide, 0.4, y, 0.08, 0.9, fill_rgb=C_RED)
    txt(slide, title, 0.65, y, 4.5, 0.35, sz=12, bold=True, color=C_RED)
    txt(slide, desc, 0.65, y + 0.32, 11.8, 0.58, sz=10.5, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE 5-PILLAR ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "The 5-Pillar Architecture", "Layered defense for enterprise-grade SAP SQL generation")

pillars = [
    (C_GREEN,  "P1", "Role-Aware RAG",
     ["AuthContext enforcement at SQL level", "Column masking & row filtering", "Immutable audit ledger", "SOX-ready logging"]),
    (C_PURPLE, "P2", "Agentic RAG",
     ["ReAct orchestration loop", "Intent classification → tool routing", "Self-correction on SQL errors", "Domain agent routing"]),
    (C_ORANGE, "P3", "Schema RAG",
     ["ChromaDB DDIC metadata retrieval", "Domain-filtered vector search", "CDS View abstraction", "Auto-trigger on empty result"]),
    (C_RED,    "P4", "SQL RAG",
     ["68+ proven SAP SQL patterns", "MANDT/LOEVM idiom injection", "Performance-aware query library", "Hot-reloadable patterns"]),
    (C_TEAL,   "P5", "Graph RAG",
     ["NetworkX cross-module traversal", "14 meta-paths (vendor, O2C, P2P)", "Node2Vec structural embeddings", "Temporal validity filters"]),
]

col_w = 2.48
for i, (col, num, name, bullets) in enumerate(pillars):
    x = 0.3 + i * col_w
    add_rect(slide, x, 1.3, col_w - 0.1, 5.9, fill_rgb=C_WHITE,
             line_rgb=RGBColor(0xD0, 0xDA, 0xE4), lw=0.5)
    add_rect(slide, x, 1.3, col_w - 0.1, 0.7, fill_rgb=col)
    txt(slide, num, x + 0.1, 1.35, 0.5, 0.3, sz=9, bold=True, color=C_WHITE)
    txt(slide, name, x + 0.1, 1.6, col_w - 0.3, 0.35, sz=12, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        y = 2.15 + j * 1.2
        add_rect(slide, x + 0.12, y, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.25, y - 0.04, col_w - 0.4, 1.0, sz=9.5, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — PILLAR 1: ROLE-AWARE SECURITY (TECHNICAL)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Pillar 1: Role-Aware RAG — 3-Layer Security Mesh",
       "Every query intercepted, validated, and audited before execution")

layers = [
    (C_GREEN, "Layer 1: Prompt Injection",
     "AuthContext hardcoded into LLM system prompt.\n"
     "BUKRS/WERKS/EKORG allowed values + masked columns.\n"
     "LLM cannot 'forget' — structurally enforced at prompt level."),
    (C_ACCENT_BLUE, "Layer 2: SQL Validation",
     "sqlglot parses generated SQL pre-execution.\n"
     "Verifies WHERE clauses exist for restricted fields.\n"
     "Blocks INSERT/UPDATE/DELETE and restricted columns."),
    (C_ACCENT_GOLD, "Layer 3: Response Masking",
     "Post-execution field redaction.\n"
     "Bank accounts, SSNs, salary data stripped.\n"
     "Final markdown sanitized before user display."),
]

for i, (col, title, desc) in enumerate(layers):
    x = 0.35 + i * 4.25
    add_rect(slide, x, 1.35, 4.0, 1.2, fill_rgb=col)
    txt(slide, title, x + 0.15, 1.4, 3.7, 0.35, sz=12, bold=True, color=C_WHITE)
    txt(slide, desc, x + 0.15, 1.78, 3.7, 0.72, sz=9.5, color=C_WHITE, italic=True)

auth_data = [
    ["Domain", "SAP Auth Object", "Maps to Column", "Filter Applied"],
    ["Finance", "F_BKPF_BUK", "BUKRS", "WHERE BUKRS IN (...)"],
    ["Material", "M_MATE_WRK", "WERKS", "WHERE WERKS IN (...)"],
    ["Purchasing", "M_BEST_EKO", "EKORG", "WHERE EKORG IN (...)"],
    ["Sales", "V_VBAK_VKO", "VKORG", "WHERE VKORG IN (...)"],
    ["BP", "F_LFA1_BUK / F_KNA1_BUK", "BUKRS", "WHERE BUKRS IN (...)"],
]
tbl = slide.shapes.add_table(6, 4, 0.35, 2.75, 12.6, 2.2).table
for ci, cw in enumerate([3.5, 3.2, 3.0, 2.9]):
    tbl.columns[ci].width = Inches(cw)
for ri, row in enumerate(auth_data):
    for ci, val in enumerate(row):
        tbl.cell(ri, ci).text = val
table_style(tbl, C_GREEN, C_LIGHT_GRAY)
txt(slide, "SAP Authorization Object → SQL WHERE Filter Mapping",
    0.35, 2.62, 12.6, 0.25, sz=10, bold=True, color=C_TEXT_DARK)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — PILLAR 2: AGENTIC ORCHESTRATOR (TECHNICAL)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Pillar 2: Agentic RAG — ReAct Orchestration Loop",
       "The autonomous brain: plan → retrieve → validate → self-correct → execute")

flow_steps = [
    (C_PURPLE, "USER\nQUERY",    "Intent Classification\n& Decomposition"),
    (RGBColor(0x4A,0x35,0x9E), "AUTH\nCHECK",  "get_auth_context()\nRole-Aware Pillar 1"),
    (C_ORANGE, "SCHEMA\nRAG",    "search_schema()\nPillar 3 — ChromaDB"),
    (C_RED,    "SQL\nRAG",       "get_sql_patterns()\nPillar 4 — Few-shot"),
    (C_TEAL,   "GRAPH\nRAG",     "traverse_graph()\nPillar 5 — NetworkX"),
    (C_ACCENT_BLUE, "SQL\nGEN",  "LLM → HANA SQL\n+ Auth Filters"),
    (C_GREEN,  "EXECUTE\n& VALIDATE","execute_hana_sql()\nSelf-Correction Loop"),
    (C_ACCENT_GOLD, "FINAL\nANSWER","Natural Language\nResponse to User"),
]

box_w = 1.45; box_h = 1.1
for i, (col, label, sub) in enumerate(flow_steps):
    row = i // 4; col_idx = i % 4
    x = 0.35 + col_idx * (box_w + 0.12)
    y = 1.4 + row * (box_h + 0.55)
    if i < 7:
        add_rect(slide, x + box_w, y + box_h / 2 - Pt(2), 0.12, Pt(4), fill_rgb=C_MID_GRAY)
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1 = tf.paragraphs[0].add_run(); r1.text = label; r1.font.size = Pt(10)
    r1.font.bold = True; r1.font.color.rgb = C_WHITE
    r2 = tf.paragraphs[0].add_run(); r2.text = f"\n{sub}"
    r2.font.size = Pt(8); r2.font.color.rgb = C_WHITE; r2.font.italic = True

tools = [
    ("get_auth_context()",   "Fetches user BUKRS/WERKS/EKORG scope from SAP Auth Profile"),
    ("search_schema()",      "ChromaDB retrieval — DDIC metadata, tables, columns, descriptions"),
    ("get_sql_patterns()",    "Few-shot proven SQL templates — MANDT, LOEVM, join idioms"),
    ("traverse_graph()",     "NetworkX shortest-path between SAP tables across modules"),
    ("execute_hana_sql()",   "Runs validated SQL on HANA, returns results or error for retry"),
    ("schema_auto_discover()","Fires when ChromaDB returns empty — reads DD08L for FK relationships"),
]
y_start = 4.15
txt(slide, "Agent Tool Ecosystem (34 tools registered)", 0.35, y_start, 12.6, 0.3, sz=11, bold=True, color=C_TEXT_DARK)
for i, (tool, desc) in enumerate(tools):
    y = y_start + 0.35 + i * 0.5
    add_rect(slide, 0.35, y + 0.06, Pt(5), Pt(5), fill_rgb=PILLAR_COLORS[i // 2])
    txt(slide, tool, 0.5, y, 2.5, 0.42, sz=10, bold=True, color=C_ACCENT_BLUE)
    txt(slide, desc, 3.0, y, 9.9, 0.42, sz=10, color=C_TEXT_DARK)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — PILLAR 3: SCHEMA RAG
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Pillar 3: Schema RAG — DDIC to ChromaDB Pipeline",
       "From 90,000 SAP tables to lean, retrieved context windows")

pipeline = [
    (C_ACCENT_BLUE, "SAP DDIC\nExtraction", "DD02L / DD03L / DD08L\nPyRFC / OData"),
    (RGBColor(0x00,0x5F,0xA3), "Document\nEnrichment", "JSON schema docs:\ntables + columns + join paths"),
    (RGBColor(0x00,0x48,0x88), "Embedding\n& Index", "all-MiniLM-L6-v2\n→ ChromaDB collection\nsap_master_schemas"),
    (C_ORANGE, "Runtime\nRetrieval", "search_schema(query,\nmodule, top_k=3)"),
]
for i, (col, label, sub) in enumerate(pipeline):
    x = 0.4 + i * 3.18
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.5, 2.9, 1.4)
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1 = tf.paragraphs[0].add_run(); r1.text = label; r1.font.size = Pt(12); r1.font.bold = True
    r1.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = f"\n{sub}"; r2.font.size = Pt(9.5); r2.font.color.rgb = C_WHITE
    if i < 3:
        add_rect(slide, x + 2.9, 2.1, 0.28, Pt(4), fill_rgb=C_MID_GRAY)

domains = [
    ["Domain", "Core Tables", "Vector Filter"],
    ["Material (MM)", "MARA, MARC, MARD, MBEW, MVKE", "module=MM"],
    ["Business Partner", "BUT000, LFA1, KNA1, ADRC, BUT100", "module=BP"],
    ["Purchasing", "EINA, EINE, EORD, EKKO, EKPO", "module=MM-PUR"],
    ["Sales & Distribution", "VBAK, VBAP, KONV, KNVV, TVRO", "module=SD"],
    ["Finance / CO", "SKA1, SKB1, CSKS, CEPC, ANLA", "module=FI"],
    ["Quality Mgmt", "QMAT, QPAC, PLMK, MAPL", "module=QM"],
    ["Project System", "PROJ, PRPS, AFVC", "module=PS"],
]
tbl2 = slide.shapes.add_table(8, 3, 0.35, 3.1, 12.6, 3.6).table
for ci, cw in enumerate([3.2, 6.0, 3.4]):
    tbl2.columns[ci].width = Inches(cw)
for ri, row in enumerate(domains):
    for ci, val in enumerate(row):
        tbl2.cell(ri, ci).text = val
table_style(tbl2, C_ORANGE, C_LIGHT_GRAY)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — PILLAR 4: SQL PATTERN RAG
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Pillar 4: SQL RAG — Few-Shot Pattern Library",
       "Grounding LLMs in proven SAP HANA SQL idioms instead of generating from scratch")

reasons = [
    "MANDT (Client) filter — LLM hallucination target #1",
    "Module-specific deletion flags — LOEVM, LVORM, LVOMA",
    "Header/Item join patterns — VBAK→VBAP→KONV",
    "Index-aware join order — BUKRS first on BSEG/MSEG",
    "Business logic — ADD_DAYS(ZFBDT, ZBD3T) for overdue dates",
]
add_rect(slide, 0.35, 1.35, 5.9, 5.8, fill_rgb=C_WHITE,
         line_rgb=RGBColor(0xD0, 0xDA, 0xE4), lw=0.5)
add_rect(slide, 0.35, 1.35, 5.9, 0.45, fill_rgb=C_RED)
txt(slide, "Critical SAP Idioms LLMs Miss", 0.5, 1.38, 5.6, 0.38, sz=11, bold=True, color=C_WHITE)
for i, r in enumerate(reasons):
    y = 1.95 + i * 0.9
    add_rect(slide, 0.5, y + 0.05, Pt(5), Pt(5), fill_rgb=C_RED)
    txt(slide, r, 0.7, y, 5.3, 0.8, sz=10.5, color=C_TEXT_DARK)

sql_lib = [
    ["Query ID", "Intent / Business Question", "Tables", "Domain"],
    ["MAT-001", "Show material base details and weight", "MARA, MAKT", "MM"],
    ["MAT-003", "Current stock for material in plant", "MARD, MARC", "MM"],
    ["SLS-001", "Sales order header and status", "VBAK, VBUK", "SD"],
    ["SLS-002", "Items and pricing for sales order", "VBAP, KONV", "SD"],
    ["FIN-001", "G/L account details for company", "SKA1, SKB1", "FI"],
    ["FIN-002", "Open vendor items (overdue)", "BSIK, LFA1", "FI"],
    ["PUR-001", "Open POs for a specific vendor", "EKKO, EKPO, LFA1", "MM"],
    ["BP-001", "Vendor bank and payment terms", "LFA1, LFB1, LFBK", "BP"],
]
tbl3 = slide.shapes.add_table(9, 4, 6.45, 1.35, 6.55, 5.8).table
for ci, cw in enumerate([0.9, 2.8, 1.7, 1.15]):
    tbl3.columns[ci].width = Inches(cw)
for ri, row in enumerate(sql_lib):
    for ci, val in enumerate(row):
        tbl3.cell(ri, ci).text = val
table_style(tbl3, C_RED, C_LIGHT_GRAY)
txt(slide, "68+ SQL Patterns Seeded Across 18 Domains", 6.45, 7.05, 6.55, 0.3, sz=9, italic=True, color=C_MID_GRAY)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — PILLAR 5: GRAPH RAG
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Pillar 5: Graph RAG — NetworkX Multi-Hop Traversal",
       "Finding join paths that Schema RAG cannot reach")

nodes = [
    (C_TEAL,    "LFA1",    "Vendor Master",     0.6,  4.3),
    (RGBColor(0x00,0x78,0x6C), "EINA",  "Purchasing\nInfo Record", 2.2, 3.5),
    (RGBColor(0x00,0x60,0x55), "MARA",  "Material Master", 3.8, 4.3),
    (RGBColor(0x00,0x48,0x40), "CSKS",  "Cost Center\nMaster", 5.4, 3.0),
]
for col, name, desc, nx, ny in nodes:
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, 1.4, 0.95)
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1 = tf.paragraphs[0].add_run(); r1.text = name; r1.font.size = Pt(12); r1.font.bold = True
    r1.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(8); r2.font.color.rgb = C_WHITE

edge_data = [
    (1.9, 4.7, 0.4, Pt(3), "LIFNR", 1.4, 4.95),
    (3.5, 4.7, 0.4, Pt(3), "MATNR", 4.15, 4.95),
    (2.95, 3.55, 0.3, Pt(3), "MATNR", 3.55, 3.85),
    (4.65, 3.95, 0.3, Pt(3), "KOSGR", 5.05, 3.45),
]
for ex, ey, ew, eh, lbl, lx, ly in edge_data:
    add_rect(slide, ex, ey, ew, eh, fill_rgb=C_TEAL)
    txt(slide, lbl, lx, ly, 0.6, 0.25, sz=8, bold=True, italic=True, color=C_TEAL, align=PP_ALIGN.CENTER)

add_rect(slide, 0.35, 1.35, 5.9, 5.8, fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0, 0xDA, 0xE4), lw=0.5)
add_rect(slide, 0.35, 1.35, 5.9, 0.45, fill_rgb=C_TEAL)
txt(slide, "NetworkX Graph — Cross-Module Edges", 0.5, 1.38, 5.6, 0.38, sz=11, bold=True, color=C_WHITE)

paths = [
    ("Procurement", "Material → Vendor",      "MARA → EINA → LFA1"),
    ("Sales",       "Material → Customer",     "MARA → KNMT → KNA1"),
    ("Costing",     "Material → Cost Center", "MARA → MARC → MBEW → CSKS"),
    ("Quality",     "Material → Inspection",   "MARA → MAPL → PLMK"),
    ("Finance",     "Vendor → G/L Account",   "LFA1 → BSIK → HKONT → SKA1"),
    ("Project",     "Material → WBS",          "MARA → RESB → PRPS"),
]
add_rect(slide, 6.45, 1.35, 6.55, 5.8, fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0, 0xDA, 0xE4), lw=0.5)
add_rect(slide, 6.45, 1.35, 6.55, 0.45, fill_rgb=C_TEAL)
txt(slide, "14 Meta-Paths — Precomputed JOIN Templates", 6.6, 1.38, 6.2, 0.38, sz=11, bold=True, color=C_WHITE)
for i, (integ, query, path) in enumerate(paths):
    y = 1.95 + i * 0.85
    add_rect(slide, 6.55, y + 0.05, Pt(5), Pt(5), fill_rgb=C_TEAL)
    txt(slide, integ, 6.75, y, 1.8, 0.35, sz=10, bold=True, color=C_TEAL)
    txt(slide, query, 6.75, y + 0.3, 2.5, 0.35, sz=9.5, italic=True, color=C_TEXT_DARK)
    txt(slide, path, 9.3, y, 3.5, 0.7, sz=9, color=C_MID_GRAY, italic=True)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — PHASE 7: TEMPORAL ENGINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Phase 7: Temporal Engine — Longitudinal Corporate Memory",
       "SAP data from 2005 becomes a living, queryable brain")

engines = [
    (C_ACCENT_BLUE, "FiscalYearEngine",
     ["4 calendar variants (4-4-5, 4-5-4, etc.)",
      "Multi-FY parsing (FY2020–FY2024)",
      "FY comparison SQL (YoY, 3-year CAGR)",
      "SAP posting date → fiscal period mapping"]),
    (C_PURPLE, "TimeSeriesAggregator",
     ["MONTHLY / QUARTERLY / YEARLY granularity",
      "Rolling 3-period averages",
      "Period-over-period delta (absolute + %)",
      "Trend detection for spend anomalies"]),
    (C_GREEN, "SupplierPerformanceIndex",
     ["Delivery reliability: EKKO/EKET/MSEG",
      "Quality: QALS UD codes (0=accepted)",
      "Price competitiveness: EINA/EINE history",
      "Composite 0-100 score per vendor"]),
    (C_TEAL, "CustomerLifetimeValueEngine",
     ["Revenue: VBAP-NETWR over 3yr window",
      "Discount sensitivity: KONV KB01/K007",
      "Payment behavior: BSID vs KNVV-ZAHLS",
      "Churn signal: 90-day order gap → flag"]),
]

for i, (col, name, bullets) in enumerate(engines):
    row = i // 2; ci = i % 2
    x = 0.35 + ci * 6.5; y = 1.35 + row * 2.9
    add_rect(slide, x, y, 6.3, 2.75, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, x, y, 6.3, 0.45, fill_rgb=col)
    txt(slide, name, x + 0.15, y + 0.05, 6.0, 0.38, sz=11, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = y + 0.55 + j * 0.52
        add_rect(slide, x + 0.15, yy + 0.05, Pt(4), Pt(4), fill_rgb=col)
        txt(slide, b, x + 0.3, yy, 5.85, 0.45, sz=9.5, color=C_TEXT_DARK)

txt(slide, "8 Historical Events Tagged: 2008 Crisis · 2011 Euro · 2015 China · COVID-19 · 2021 Supply Chain · 2022 Inflation",
    0.35, 7.1, 12.6, 0.3, sz=9, italic=True, color=C_MID_GRAY)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — PHASE 8: QM + NEGOTIATION + CONFIDENCE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Phase 8: QM Semantic + Negotiation + Confidence Scoring",
       "Domain intelligence that generic agents cannot match out-of-the-box")

left_col = [
    (C_ORANGE, "QM Semantic Search",
     ["20-year QM notification long-text embeddings",
      "ChromaDB collection: chroma_qm_db",
      "QMEL-LONGTEXT semantic similarity",
      "Fires on: inspection, defect, complaint keywords"]),
    (C_ACCENT_GOLD, "Negotiation Briefing",
     ["CLV Tier: GOLD / SILVER / BRONZE",
      "PSI Score, Churn Risk, BATNA analysis",
      "Target price increase + acceptance floor",
      "Top tactic selection from 8-strategy library"]),
]

for i, (col, name, bullets) in enumerate(left_col):
    y = 1.35 + i * 2.6
    add_rect(slide, 0.35, y, 6.3, 2.45, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, 0.35, y, 6.3, 0.45, fill_rgb=col)
    txt(slide, name, 0.5, y + 0.05, 6.0, 0.38, sz=11, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = y + 0.55 + j * 0.45
        add_rect(slide, 0.5, yy + 0.05, Pt(4), Pt(4), fill_rgb=col)
        txt(slide, b, 0.65, yy, 5.85, 0.4, sz=9.5, color=C_TEXT_DARK)

# Confidence Score breakdown
add_rect(slide, 6.85, 1.35, 6.1, 5.8, fill_rgb=C_WHITE, line_rgb=C_ACCENT_BLUE, lw=1)
add_rect(slide, 6.85, 1.35, 6.1, 0.45, fill_rgb=C_ACCENT_BLUE)
txt(slide, "6-Signal Confidence Score (Fully Auditable)", 7.0, 1.38, 5.8, 0.38, sz=11, bold=True, color=C_WHITE)

signals = [
    ("critique",       "30%", "SQL gatekeeper — 7-point check"),
    ("result_density", "25%", "Rows returned vs. table size"),
    ("routing_path",   "15%", "Meta-path fast path vs. full RAG"),
    ("self_heal",      "10%", "Self-healer applied (Y/N)"),
    ("temporal",       "10%", "Temporal engine triggered"),
    ("schema_breadth", "10%", "Number of tables involved"),
]

for i, (sig, wt, desc) in enumerate(signals):
    y = 1.95 + i * 0.8
    add_rect(slide, 7.0, y, 0.9, 0.35, fill_rgb=C_ACCENT_BLUE)
    txt(slide, sig, 7.0, y, 0.9, 0.35, sz=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, f"wt={wt}", 7.95, y, 0.7, 0.35, sz=9, bold=True, color=C_ACCENT_BLUE)
    txt(slide, desc, 8.7, y, 4.1, 0.35, sz=9.5, color=C_TEXT_DARK)

txt(slide, "Example: composite=0.83 | critique:0.86 | density:0.91 | routing:1.00",
    6.85, 6.85, 6.1, 0.3, sz=9, italic=True, color=C_MID_GRAY)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — SAP JOULE VS. OUR BUILD (TECHNICAL)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "SAP Joule vs. Know Your SAP Masters — Technical Comparison",
       "Where our build matches, exceeds, or complements SAP native")

comp = [
    ["Dimension", "SAP Joule / Business AI", "Know Your SAP Masters", "Advantage"],
    ["Auth Model", "SAP AuthContext (RFC-based)", "SAPAuthContext class (same concept)", "Parity"],
    ["Column Masking", "AuthObject-based", "result_mask() per role", "Parity"],
    ["Tool Registry", "2,100+ AI Skills", "34+ tools in TOOL_REGISTRY", "SAP (scale)"],
    ["Self-Correction", "Implicit 'reflection'", "10 explicit SQL repair rules", "Us (auditable)"],
    ["Schema Retrieval", "SAP Knowledge Graph (proprietary)", "ChromaDB — swap-in-place", "Parity"],
    ["SQL Patterns", "Workflow-level skills", "Executable HANA SQL patterns", "Us (SQL layer)"],
    ["Graph Traversal", "Enterprise graph (generic)", "14 meta-paths + Node2Vec", "Us (domain depth)"],
    ["Temporal Engine", "Date-effective tables (SAP native)", "Temporal engine + SQL filter gen", "Parity"],
    ["Confidence Score", "Black-box 'reflection'", "6-signal explicit breakdown", "Us (auditable)"],
    ["CDS Views", "Native S/4HANA Cloud", "Phase roadmap item", "SAP (cloud-ready)"],
]

tbl = slide.shapes.add_table(len(comp), 4, 0.35, 1.35, 12.6, 5.5).table
for ci, cw in enumerate([2.5, 4.2, 4.2, 1.7]):
    tbl.columns[ci].width = Inches(cw)
for ri, row in enumerate(comp):
    for ci, val in enumerate(row):
        tbl.cell(ri, ci).text = val
table_style(tbl, C_NAVY, C_LIGHT_GRAY)

for ri, row in enumerate(comp):
    if ri == 0:
        continue
    cell = tbl.cell(ri, 3)
    cell.fill.solid()
    advantage = row[3]
    if advantage == "Parity":
        cell.fill.fore_color.rgb = C_MID_GRAY
    elif advantage == "SAP (scale)" or advantage == "SAP (cloud-ready)":
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
    else:
        cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xD5)
    for para in cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(9); run.font.bold = True
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — SELF-HEALING + ERROR REPAIR RULES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Self-Healer — 10 Autonomous SQL Repair Rules",
       "Fires at three trigger points: critique FAIL · validation ERROR · execution ERROR")

rules = [
    ("Rule 1", "Missing MANDT",
     "If MANDT not in WHERE → auto-inject MANDT='800'"),
    ("Rule 2", "Missing LOEVM_K",
     "If vendor table and no LOEVM_K → append AND LOEVM_K = ''"),
    ("Rule 3", "Wrong date format",
     "If date literal not YYYYMMDD → convert using DATS_CAST"),
    ("Rule 4", "Missing index column in JOIN",
     "If BUKRS not first in BSEG join → reorder to BUKRS first"),
    ("Rule 5", "Duplicate column in SELECT",
     "sqlglot AST parse → deduplicate SELECT list"),
    ("Rule 6", "N+1 subquery detected",
     "If correlated subquery on MSEG → rewrite as JOIN"),
    ("Rule 7", "Missing WHERE on restricted field",
     "If BUKRS not in WHERE → inject AuthContext BUKRS filter"),
    ("Rule 8", "Syntax error — quote mismatch",
     "If single-quote orphan → escape or wrap with REGEX"),
    ("Rule 9", "Execution timeout",
     "If query > 30s → inject HINT (NO_PARALLEL) and rerun"),
    ("Rule 10", "Empty result — wrong table",
     "If 0 rows and schema confidence < 0.5 → trigger schema_auto_discover()"),
]

for i, (num, title, fix) in enumerate(rules):
    row = i // 2; ci = i % 2
    x = 0.35 + ci * 6.5; y = 1.35 + row * 1.15
    add_rect(slide, x, y, 6.3, 1.05, fill_rgb=C_WHITE, line_rgb=C_RED, lw=0.5)
    add_rect(slide, x, y, 0.45, 1.05, fill_rgb=C_RED)
    txt(slide, num, x, y, 0.45, 1.05, sz=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, x + 0.55, y + 0.05, 5.6, 0.3, sz=10, bold=True, color=C_RED)
    txt(slide, fix, x + 0.55, y + 0.38, 5.6, 0.6, sz=9.5, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — IMPLEMENTATION PHASES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Implementation Roadmap", "Build fast, add cross-module intelligence incrementally")

phases = [
    (C_ACCENT_BLUE, "Phase 1", "Weeks 1–4",
     "Isolated Master Data Foundations",
     ["Schema RAG + SQL RAG for BP, Material, Finance",
      "Role-Aware RAG core — AuthContext service",
      "Agentic Orchestrator — basic ReAct loop",
      "Fast win: Chat with single-domain tables"]),
    (RGBColor(0x00, 0x7A, 0xC2), "Phase 2", "Weeks 5–8",
     "Cross-Module via Graph RAG",
     ["NetworkX graph construction from DD08L",
      "Graph RAG for MM→BP→FI join paths",
      "Multi-hop Q&A: MARA→EINA→LFA1→CSKS",
      "CDS View abstraction layer"]),
    (RGBColor(0x00, 0x60, 0x99), "Phase 3", "Weeks 9–11",
     "Transactional Data Integration",
     ["Extend to PO, Invoice, Material Docs, JEs",
      "High-volume table support: EKKO/EKPO, BSEG",
      "SQL Library expansion — aggregation patterns",
      "Response masking refinement"]),
    (C_DARK_NAVY, "Phase 4", "Weeks 12–14",
     "Production Polish & Audit",
     ["Immutable audit ledger — SOX ready",
      "UAT with distinct persona roles",
      "Performance benchmarking & optimization",
      "Production deployment & monitoring"]),
]

for i, (col, phase, timing, title, bullets) in enumerate(phases):
    x = 0.3 + i * 3.25
    add_rect(slide, x, 1.35, 3.1, 5.8, fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0, 0xDA, 0xE4), lw=0.5)
    add_rect(slide, x, 1.35, 3.1, 1.05, fill_rgb=col)
    txt(slide, phase, x + 0.1, 1.38, 1.4, 0.3, sz=9, bold=True, color=C_WHITE)
    txt(slide, timing, x + 1.5, 1.38, 1.4, 0.3, sz=9, color=C_ACCENT_GOLD)
    txt(slide, title, x + 0.1, 1.68, 2.9, 0.6, sz=11, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        y = 2.55 + j * 1.05
        add_rect(slide, x + 0.12, y + 0.06, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.25, y, 2.75, 0.95, sz=9.5, color=C_TEXT_DARK)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING CTA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
add_rect(slide, 9.2, 0, 4.13, H, fill_rgb=C_NAVY)
add_rect(slide, 9.13, 0, Pt(6), H, fill_rgb=C_ACCENT_BLUE)
add_rect(slide, 0.55, 1.6, 3.5, Pt(3), fill_rgb=C_ACCENT_GOLD)

txt(slide, "Built for Basis Teams\nWho Demand Auditability",
    0.5, 1.75, 8.4, 1.5, sz=34, bold=True, color=C_WHITE)
txt(slide,
    "5-Pillar RAG gives your Basis team SQL-level visibility, role-level security enforcement, "
    "and a fully auditable confidence score on every query.",
    0.5, 3.4, 8.4, 1.0, sz=14, color=C_LIGHT_BLUE, italic=True)

stats = [
    ("34+", "Tools in\nRegistry"),
    ("68+", "SQL Patterns\nSeeded"),
    ("14", "Meta-Paths\nDefined"),
    ("6-Signal", "Confidence\nBreakdown"),
]
for i, (num, lbl) in enumerate(stats):
    x = 0.5 + i * 2.1
    txt(slide, num, x, 4.6, 1.8, 0.9, sz=36, bold=True, color=C_ACCENT_BLUE, align=PP_ALIGN.CENTER)
    txt(slide, lbl, x, 5.45, 1.8, 0.6, sz=10, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Immediate Next Steps", 9.45, 1.5, 3.5, 0.5, sz=18, bold=True, color=C_WHITE)
next_steps = [
    "Deploy Schema RAG for vendor master (LFA1 domain)",
    "Connect real SAP HANA via PyRFC / hdbcli",
    "Seed first 10 validated SQL patterns",
    "Onboard AP_CLERK and BUYER roles for AuthContext",
    "Run integration test suite (pytest)",
]
for i, s in enumerate(next_steps):
    y = 2.15 + i * 0.9
    add_rect(slide, 9.45, y + 0.08, Pt(5), Pt(5), fill_rgb=C_ACCENT_BLUE)
    txt(slide, s, 9.65, y, 3.3, 0.8, sz=10.5, color=C_WHITE)
footer(slide, "Know Your SAP Masters — Basis Team Deck")


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
out = r"C:\Users\vishnu\.openclaw\workspace\SAP_HANA_LLM_VendorChatbot\KnowYourSAPMasters_BasisTeam_Deck.pptx"
prs.save(out)
print(f"[OK] Basis Team Deck saved: {out} | Slides: {len(prs.slides)}")
