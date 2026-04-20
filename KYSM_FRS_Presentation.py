"""
KYSM FRS v1.0 — NRL-Branded PowerPoint Presentation Generator
Uses NRL brand guidelines from NRLlogoBooklet/SKILL.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy
from datetime import date

# ── NRL Brand Colors ──────────────────────────────────────────────────────────
NRL_YELLOW = RGBColor(0xFF, 0xD1, 0x00)   # #FFD100 — Sun
NRL_RED    = RGBColor(0xFF, 0x4A, 0x00)   # #FF4A00 — Blazing trail
NRL_BLUE   = RGBColor(0x00, 0x3D, 0xA5)   # #003DA5 — Dark navy blue
NRL_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NRL_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
NRL_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
NRL_DNAVY  = RGBColor(0x00, 0x2D, 0x72)   # Dark Navy headers
NRL_ACCENT = RGBColor(0x00, 0x2D, 0x72)   # same as dark navy for accents

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_name="Calibri",
                font_size=18, bold=False, italic=False, color=NRL_BLACK,
                align=PP_ALIGN.LEFT, wrap=True, v_anchor=None):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_accent_bar(slide, left, top, width, height=Pt(4)):
    """Yellow accent bar."""
    add_rect(slide, left, top, width, height, fill_color=NRL_YELLOW)

def add_slide_number(slide, num, total):
    """Add slide number bottom-right."""
    add_textbox(slide,
               SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.55),
               Inches(1.1), Inches(0.4),
               f"Slide {num} of {total}",
               font_name="Calibri", font_size=10,
               color=RGBColor(0x80, 0x80, 0x80),
               align=PP_ALIGN.RIGHT)

def add_footer(slide, left_text="Numaligarh Refinery Limited | Internal Use Only"):
    """Add footer bar at bottom of slide."""
    bar_height = Inches(0.38)
    add_rect(slide, 0, SLIDE_H - bar_height, SLIDE_W, bar_height, fill_color=NRL_DNAVY)
    add_textbox(slide,
                Inches(0.3), SLIDE_H - bar_height + Inches(0.05),
                Inches(6), Inches(0.3),
                left_text,
                font_name="Calibri", font_size=9,
                color=NRL_WHITE, align=PP_ALIGN.LEFT)

def add_title_slide_header(slide, title_text, subtitle_text=None):
    """Add a clean title block to a content slide."""
    # Top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.35), SLIDE_W - Inches(1), Inches(0.7),
                title_text, font_name="Montserrat", font_size=32, bold=True,
                color=NRL_DNAVY, align=PP_ALIGN.LEFT)
    # Accent underline
    add_accent_bar(slide, Inches(0.5), Inches(1.0), Inches(2.5))
    if subtitle_text:
        add_textbox(slide, Inches(0.5), Inches(1.1), SLIDE_W - Inches(1), Inches(0.4),
                    subtitle_text, font_name="Calibri", font_size=16,
                    color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.LEFT)

def add_bullet_block(slide, left, top, width, height, items,
                     font_name="Calibri", font_size=15, color=NRL_BLACK,
                     bullet_color=NRL_YELLOW, indent=Inches(0.25),
                     line_spacing=None):
    """Add a series of bullet items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"\u25cf  {item}"
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def add_2col_table(slide, left, top, width, height, headers, rows,
                   col_widths=None, header_fill=NRL_BLUE, row_fill=NRL_WHITE,
                   alt_fill=RGBColor(0xE8, 0xF0, 0xFB)):
    """Add a styled table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Column widths
    if col_widths:
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = NRL_WHITE

    # Data rows
    for ri, row in enumerate(rows):
        fill = alt_fill if ri % 2 == 0 else row_fill
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(11)
            run.font.color.rgb = NRL_BLACK

    return table

# ── Build Presentation ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

TOTAL_SLIDES = 20  # approximate

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, NRL_DNAVY)

# Top yellow stripe
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=NRL_YELLOW)

# Center content block
add_textbox(slide, Inches(1), Inches(1.8), SLIDE_W - Inches(2), Inches(1.2),
            "KYSM — Know Your SAP Masters",
            font_name="Montserrat", font_size=44, bold=True,
            color=NRL_WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.1), SLIDE_W - Inches(2), Inches(0.7),
            "Agentic RAG System for SAP S/4 HANA",
            font_name="Montserrat", font_size=26, bold=False,
            color=NRL_YELLOW, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.85), SLIDE_W - Inches(2), Inches(0.5),
            "Functional Requirement Specification v1.0",
            font_name="Calibri", font_size=18,
            color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(4), Inches(4.45), Inches(5.333), Inches(0.04), fill_color=NRL_YELLOW)

add_textbox(slide, Inches(1), Inches(4.6), SLIDE_W - Inches(2), Inches(0.4),
            "Numaligarh Refinery Limited",
            font_name="Calibri", font_size=14,
            color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.05), SLIDE_W - Inches(2), Inches(0.35),
            "Document ID: KYSM-FRS-001  |  Status: Consolidated",
            font_name="Calibri", font_size=12,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.4), SLIDE_W - Inches(2), Inches(0.35),
            "2026-04-16",
            font_name="Calibri", font_size=12,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# Bottom tagline bar
add_rect(slide, 0, SLIDE_H - Inches(0.7), SLIDE_W, Inches(0.7), fill_color=RGBColor(0x00, 0x1A, 0x4A))
add_textbox(slide, Inches(1), SLIDE_H - Inches(0.62), SLIDE_W - Inches(2), Inches(0.5),
            "Energy in Motion",
            font_name="Montserrat", font_size=16, bold=True,
            color=NRL_YELLOW, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "Agenda", "What's inside this document")

agenda_items = [
    ("01", "Introduction & Purpose",        "Scope, definitions, document conventions"),
    ("02", "Architecture Overview",          "System perspective, product functions"),
    ("03", "Five-Pillar RAG Architecture",   "Security, Orchestrator, Schema, SQL, Graph"),
    ("04", "Multi-Agent Swarm",              "7 domain agents, planner, synthesis"),
    ("05", "26-Phase Orchestration",         "Canonical pipeline from Meta-Path to Commit"),
    ("06", "Self-Critique / Self-Heal",     "10 error codes, 6 heal strategies"),
    ("07", "Security & Sandbox",             "AuthContext, 7-layer sandbox, Sentinel"),
    ("08", "Graph RAG & Memgraph",           "Path discovery, Node2Vec, dual-mode"),
    ("09", "Persistent Memory & Dialog",     "6 stores, session management"),
    ("10", "NFRs & Compliance",              "Performance, security, Willison scorecard"),
    ("11", "Pending Requirements & Summary","Roadmap P0–P3, next steps"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.35) + i * Inches(0.49)
    # Number badge
    badge = add_rect(slide, Inches(0.5), y, Inches(0.45), Inches(0.38), fill_color=NRL_BLUE)
    add_textbox(slide, Inches(0.5), y + Inches(0.04), Inches(0.45), Inches(0.35),
                num, font_name="Montserrat", font_size=11, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.05), y + Inches(0.03), Inches(3.5), Inches(0.35),
                title, font_name="Montserrat", font_size=14, bold=True,
                color=NRL_DNAVY, align=PP_ALIGN.LEFT)
    # Desc
    add_textbox(slide, Inches(4.6), y + Inches(0.03), Inches(8.0), Inches(0.35),
                desc, font_name="Calibri", font_size=12,
                color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.LEFT)

add_footer(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION & PURPOSE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "01  Introduction & Purpose",
                        "What KYSM is and why it exists")

# Left column — What is KYSM
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
            "What is KYSM?", font_name="Montserrat", font_size=18, bold=True,
            color=NRL_DNAVY)
add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.04), fill_color=NRL_YELLOW)

what_items = [
    "Agentic Retrieval-Augmented-Generation (RAG) system",
    "Translates natural-language questions into SAP HANA SQL",
    "Executes read-only queries with role-scoped access control",
    "Sits between upstream portal and SAP S/4 HANA backend",
    "NOT a chatbot — an autonomous SQL engineering system",
]
add_bullet_block(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.5), what_items)

# Right column — Key Stats
add_rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(3.2),
         fill_color=NRL_DNAVY)
stats = [
    ("18", "SAP Modules Covered"),
    ("26", "Runtime Phases"),
    ("7",  "Domain Agents"),
    ("52", "Total Tools"),
    ("114","DDIC Tables"),
    ("6",  "Memory Stores"),
]
for i, (val, lbl) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = Inches(7.1) + col * Inches(2.9)
    y = Inches(1.5) + row * Inches(0.95)
    add_textbox(slide, x, y, Inches(1.5), Inches(0.6),
                val, font_name="Montserrat", font_size=30, bold=True,
                color=NRL_YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.5), y + Inches(0.1), Inches(1.3), Inches(0.5),
                lbl, font_name="Calibri", font_size=11,
                color=NRL_WHITE)

# Out of scope box
add_rect(slide, Inches(0.5), Inches(4.55), Inches(5.8), Inches(1.6),
         fill_color=RGBColor(0xF0, 0xF4, 0xFF),
         line_color=NRL_BLUE, line_width=Pt(1))
add_textbox(slide, Inches(0.65), Inches(4.65), Inches(5.5), Inches(0.35),
            "Out of Scope for v1.0", font_name="Montserrat", font_size=13, bold=True,
            color=NRL_BLUE)
scope_items = [
    "Write-back via BAPI (LFA1, MARA)",
    "Real-time streaming changes",
    "Non-HANA backends",
    "End-user authentication (delegated to portal)",
]
add_bullet_block(slide, Inches(0.65), Inches(5.05), Inches(5.5), Inches(1.2), scope_items,
                 font_size=12)

add_footer(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "02  Architecture Overview",
                        "System perspective — KYSM as a backend service")

# Architecture box diagram (drawn with rectangles)
box_defs = [
    # (label, x_frac, y_frac, w_in, h_in, fill, text_color)
    ("Chat UI / API\nClient",    0.03, 0.25, 1.6, 1.0,  RGBColor(0xE8, 0xF0, 0xFB), NRL_BLUE),
    ("KYSM FastAPI\nBackend",    0.22, 0.25, 1.6, 1.0,  NRL_DNAVY,                  NRL_WHITE),
    ("Memgraph\n(graph)",        0.41, 0.10, 1.4, 0.8,  RGBColor(0xFF, 0xF0, 0xE0), NRL_RED),
    ("Qdrant\n(vectors)",        0.41, 0.50, 1.4, 0.8,  RGBColor(0xE8, 0xFF, 0xE8), RGBColor(0x00, 0x60, 0x00)),
    ("RabbitMQ\n(task queue)",   0.60, 0.10, 1.4, 0.8,  RGBColor(0xF5, 0xF5, 0xFF), NRL_BLUE),
    ("SAP HANA\n(hdbcli)",       0.60, 0.50, 1.4, 0.8,  RGBColor(0xFF, 0xFF, 0xE8), RGBColor(0x60, 0x40, 0x00)),
]

for (lbl, xf, yf, w, h, fill, tc) in box_defs:
    bx = SLIDE_W * xf
    by = SLIDE_H * yf
    add_rect(slide, bx, by, Inches(w), Inches(h), fill_color=fill,
             line_color=NRL_BLUE, line_width=Pt(0.75))
    add_textbox(slide, bx + Inches(0.05), by + Inches(0.15),
                Inches(w - 0.1), Inches(h - 0.2),
                lbl, font_name="Calibri", font_size=11, bold=True,
                color=tc, align=PP_ALIGN.CENTER)

# Arrows (horizontal lines with text labels)
arrow_pairs = [
    (0.165, 0.55, 0.215, 0.55, "HTTP"),
    (0.375, 0.55, 0.405, 0.55, "Bolt"),
    (0.565, 0.35, 0.595, 0.35, "gRPC"),
    (0.565, 0.75, 0.595, 0.75, "AMQP"),
    (0.745, 0.55, 0.775, 0.55, "hdbcli"),
]
for (x1f, y1f, x2f, y2f, lbl) in arrow_pairs:
    x1, y1 = SLIDE_W * x1f, SLIDE_H * y1f
    x2, y2 = SLIDE_W * x2f, SLIDE_H * y2f
    add_rect(slide, x1, y1 - Inches(0.02), x2 - x1, Inches(0.04),
             fill_color=NRL_YELLOW)
    mx = (x1 + x2) / 2
    add_textbox(slide, mx - Inches(0.3), y1 - Inches(0.22), Inches(0.6), Inches(0.2),
                lbl, font_name="Calibri", font_size=8, bold=True,
                color=NRL_BLACK, align=PP_ALIGN.CENTER)

# Product Functions table (below diagram)
add_textbox(slide, Inches(0.5), Inches(4.25), Inches(12.3), Inches(0.35),
            "High-Level Product Functions", font_name="Montserrat", font_size=16, bold=True,
            color=NRL_DNAVY)
add_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.04), fill_color=NRL_YELLOW)

fn_headers = ["ID", "Function", "Description"]
fn_rows = [
    ["F1", "Translate NL → SQL",       "Natural language to safe, role-scoped SAP HANA SQL"],
    ["F2", "Validate via 7-layer sandbox", "Validate before execution"],
    ["F3", "Execute & Mask Results",    "Return masked, role-appropriate data"],
    ["F4", "Multi-turn Memory",          "Per-role conversational context"],
    ["F5", "Self-Critique & Self-Heal",  "Autonomous SQL repair"],
    ["F6", "Learn Patterns",             "Promote/demote SQL patterns over time"],
    ["F7", "Threat Sentinel",            "Behavioral anomaly detection"],
    ["F8", "Multi-Agent Swarm",          "Cross-module query execution"],
]
add_2col_table(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.6),
               fn_headers, fn_rows,
               col_widths=[Inches(0.8), Inches(2.8), Inches(8.7)])

add_footer(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FIVE-PILLAR RAG ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "03  Five-Pillar RAG Architecture",
                        "Concurrent retrieval pillars + one hybrid")

pillars = [
    ("1", "Security\nPillar",        "SAPAuthContext\nIn-memory + Redis",
     "Row/table/column scope,\ndenied tables,\nmasked fields",          NRL_BLUE),
    ("2", "Orchestrator\nPillar",     "Stateful Controller",
     "Phase routing,\nmemory wiring,\nSwarm decisions",            NRL_DNAVY),
    ("3", "Schema RAG\nPillar",       "Qdrant — sap_schema",
     "Table/field retrieval,\nDDIC metadata,\nsemantic search",         RGBColor(0x00, 0x60, 0x00)),
    ("4", "SQL Pattern\nPillar",      "Qdrant — sql_patterns",
     "68+ proven templates,\nfew-shot injection,\ndomain-specific",         RGBColor(0x80, 0x00, 0x80)),
    ("5", "Graph RAG\nPillar",        "Memgraph + NetworkX",
     "Path discovery,\nmeta-paths (14),\nSteiner tree",              RGBColor(0xCC, 0x55, 0x00)),
    ("5½", "Graph Embed\nPillar",    "Qdrant — node+text",
     "Node2Vec 64-dim,\nstructural roles,\nhybrid ANN search",       RGBColor(0x00, 0x70, 0x70)),
]

col_w = Inches(2.0)
start_x = Inches(0.35)
for i, (num, title, store, desc, color) in enumerate(pillars):
    x = start_x + i * (col_w + Inches(0.06))
    # Card background
    add_rect(slide, x, Inches(1.35), col_w, Inches(4.5),
             fill_color=RGBColor(0xF8, 0xF8, 0xFF),
             line_color=color, line_width=Pt(1.5))
    # Top color band
    add_rect(slide, x, Inches(1.35), col_w, Inches(0.55), fill_color=color)
    # Number
    add_textbox(slide, x, Inches(1.37), col_w, Inches(0.5),
                f"Pillar {num}", font_name="Montserrat", font_size=13, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.1), Inches(1.95), col_w - Inches(0.2), Inches(0.7),
                title, font_name="Montserrat", font_size=13, bold=True,
                color=color, align=PP_ALIGN.CENTER)
    # Store label
    add_rect(slide, x + Inches(0.1), Inches(2.65), col_w - Inches(0.2), Inches(0.38),
             fill_color=RGBColor(0xEE, 0xF2, 0xFF))
    add_textbox(slide, x + Inches(0.1), Inches(2.67), col_w - Inches(0.2), Inches(0.35),
                store, font_name="Calibri", font_size=9, bold=True,
                color=NRL_BLUE, align=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.1), Inches(3.1), col_w - Inches(0.2), Inches(2.6),
                desc, font_name="Calibri", font_size=11,
                color=NRL_BLACK, align=PP_ALIGN.CENTER)

# Note
add_textbox(slide, Inches(0.5), Inches(6.05), SLIDE_W - Inches(1), Inches(0.35),
            "Note: ChromaDB deprecated for Schema & SQL Pattern RAG (as of M9). "
            "All vector collections now served exclusively by Qdrant.",
            font_name="Calibri", font_size=11, italic=True,
            color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.LEFT)

add_footer(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MULTI-AGENT SWARM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "04  Multi-Agent Swarm",
                        "7 domain agents with Planner + Synthesis coordination")

# Left — Agents
agents = [
    ("Supervisor",         "Routes query to correct agent(s)",    NRL_DNAVY),
    ("Planner Agent",      "Decomposes CROSS_MODULE into DAG",    NRL_BLUE),
    ("bp — Business Partner", "LFA1, KNA1, BUT000, ADRC",       RGBColor(0x00, 0x70, 0x00)),
    ("mm — Material Master",  "MARA, MARC, MARD, MBEW",         RGBColor(0xCC, 0x55, 0x00)),
    ("pur — Purchasing",      "EKKO, EKPO, EKBE, EINA",          RGBColor(0x80, 0x00, 0x80)),
    ("sd — Sales & Dist.",    "VBAK, VBAP, VBRK, LIKP",         RGBColor(0x00, 0x60, 0x80)),
    ("qm — Quality Mgmt",     "QALS, QMEL, QPAM",                RGBColor(0xCC, 0x00, 0x00)),
    ("wm — Warehouse Mgmt",   "LAGP, LTAK, LQUA",                 RGBColor(0x80, 0x60, 0x00)),
    ("cross — Cross Module", "BSEG, BKPF, COSP, COSS",          RGBColor(0x00, 0x50, 0x80)),
    ("Synthesis Agent",      "Merges results, deduplicates",     NRL_RED),
]

for i, (name, desc, color) in enumerate(agents):
    y = Inches(1.35) + i * Inches(0.53)
    add_rect(slide, Inches(0.4), y, Inches(0.18), Inches(0.42), fill_color=color)
    add_textbox(slide, Inches(0.68), y + Inches(0.03), Inches(3.5), Inches(0.38),
                name, font_name="Montserrat", font_size=12, bold=True,
                color=NRL_BLACK, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(4.2), y + Inches(0.03), Inches(3.0), Inches(0.38),
                desc, font_name="Calibri", font_size=11,
                color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.LEFT)

# Right — SwarmDecision table
add_textbox(slide, Inches(7.5), Inches(1.35), Inches(5.5), Inches(0.35),
            "SwarmDecision Routing", font_name="Montserrat", font_size=15, bold=True,
            color=NRL_DNAVY)
add_rect(slide, Inches(7.5), Inches(1.68), Inches(5.5), Inches(0.04), fill_color=NRL_YELLOW)

dec_headers = ["Decision",      "Trigger",                      "Execution"]
dec_rows = [
    ["SINGLE",        "1 agent confidence ≥ 0.7",           "That agent only"],
    ["PARALLEL",      "≥2 domain signals + complexity kw",  "ThreadPoolExecutor"],
    ["CROSS_MODULE",  "≥2 entity types or bridge keywords", "Full Graph RAG"],
    ["FALLBACK",      "No agent ≥ 0.4 confidence",          "Standard orchestrator"],
]
add_2col_table(slide, Inches(7.5), Inches(1.78), Inches(5.5), Inches(2.2),
               dec_headers, dec_rows,
               col_widths=[Inches(1.3), Inches(2.6), Inches(1.6)])

# Authority weights
add_textbox(slide, Inches(7.5), Inches(4.15), Inches(5.5), Inches(0.35),
            "Agent Authority Weights", font_name="Montserrat", font_size=15, bold=True,
            color=NRL_DNAVY)
add_rect(slide, Inches(7.5), Inches(4.48), Inches(5.5), Inches(0.04), fill_color=NRL_YELLOW)

auth_rows = [
    ["Planner",       "8",  "Highest — owns routing DAG"],
    ["Domain Specialists", "7", "Outrank cross_agent on their data"],
    ["Synthesis Agent", "6", "Merges and ranks cross-domain results"],
    ["cross_agent",   "5",  "Cannot outrank specialists"],
]
add_2col_table(slide, Inches(7.5), Inches(4.58), Inches(5.5), Inches(1.6),
               ["Agent", "Weight", "Role"], auth_rows,
               col_widths=[Inches(1.8), Inches(0.7), Inches(3.0)])

add_footer(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 26-PHASE ORCHESTRATION PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "05  26-Phase Orchestration Pipeline",
                        "Canonical runtime phases — every query traverses all in order")

phases_left = [
    ("0",  "Meta-Path Match",        "Fast-path template lookup"),
    ("1",  "Schema RAG",             "Qdrant table retrieval"),
    ("1b", "DDIC Auto-Discover",     "Fallback when Phase 1 empty"),
    ("1.5","Graph Embedding Search","Node2Vec + text hybrid ANN"),
    ("1.75","QM Semantic Search",   "Quality-specific semantic path"),
    ("2",  "SQL Pattern RAG",        "Proven template retrieval"),
    ("2b", "Temporal Detection",     "Date / fiscal-Q detection"),
    ("2c", "Temporal Analysis",      "Filter construction"),
    ("2d", "Negotiation Briefing",   "Pre-swarm context pack"),
    ("3",  "Graph RAG",              "AllPathsExplorer / Steiner tree"),
    ("4",  "SQL Assembly",           "Inject AuthContext + filters"),
    ("4.5","Self-Critique",          "7-point gate (score ≥ 0.7)"),
    ("5",  "SQL Validation",         "DML/MANDT/AuthContext gate"),
]

phases_right = [
    ("5.5","Dry-Run Harness",       "SELECT COUNT(*) trial"),
    ("6",  "Execution",             "Route to HANA via proxy"),
    ("6a", "Self-Heal Retry",       "10-rule repair loop"),
    ("6b", "Result Fetch",          "Stream or page rows"),
    ("7",  "Result Masking",         "AuthContext column redaction"),
    ("8",  "Synthesis",             "Natural-language answer"),
    ("8b", "Memory Log",             "Query + pattern + gotcha persistence"),
    ("9",  "Self-Improver Review",  "Promote / demote / ghost-inject"),
    ("10", "File-Based Handoff",    "Context isolation for next turn"),
    ("11", "Trajectory Log",         "Per-query replay record"),
    ("12", "Quality Metrics Eval",  "Critique + masking + heal stats"),
    ("12b","Meta-Harness Loop",     "Offline retraining hooks"),
    ("13","Agent Bus Publish",      "Broadcast to swarm (if swarm)"),
    ("13b","Negotiation-Commit",   "Final commit after negotiation"),
]

for i, (num, name, desc) in enumerate(phases_left):
    y = Inches(1.3) + i * Inches(0.42)
    badge_col = NRL_BLUE if i % 2 == 0 else NRL_DNAVY
    add_rect(slide, Inches(0.4), y, Inches(0.55), Inches(0.35), fill_color=badge_col)
    add_textbox(slide, Inches(0.4), y + Inches(0.04), Inches(0.55), Inches(0.3),
                num, font_name="Montserrat", font_size=10, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.05), y + Inches(0.04), Inches(2.5), Inches(0.32),
                name, font_name="Montserrat", font_size=11, bold=True,
                color=NRL_DNAVY)
    add_textbox(slide, Inches(3.6), y + Inches(0.04), Inches(3.0), Inches(0.32),
                desc, font_name="Calibri", font_size=10,
                color=RGBColor(0x44, 0x44, 0x44))

for i, (num, name, desc) in enumerate(phases_right):
    y = Inches(1.3) + i * Inches(0.42)
    badge_col = NRL_RED if i % 2 == 0 else RGBColor(0xCC, 0x55, 0x00)
    add_rect(slide, Inches(6.8), y, Inches(0.55), Inches(0.35), fill_color=badge_col)
    add_textbox(slide, Inches(6.8), y + Inches(0.04), Inches(0.55), Inches(0.3),
                num, font_name="Montserrat", font_size=10, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.45), y + Inches(0.04), Inches(2.5), Inches(0.32),
                name, font_name="Montserrat", font_size=11, bold=True,
                color=NRL_DNAVY)
    add_textbox(slide, Inches(10.0), y + Inches(0.04), Inches(3.0), Inches(0.32),
                desc, font_name="Calibri", font_size=10,
                color=RGBColor(0x44, 0x44, 0x44))

# Centre divider
add_rect(slide, Inches(6.65), Inches(1.3), Inches(0.04), Inches(6.0), fill_color=NRL_YELLOW)

# Note
add_textbox(slide, Inches(0.5), Inches(7.0), SLIDE_W - Inches(1), Inches(0.35),
            "Phase 4.5 Self-Critique 7-point gate: (1) SELECT-only  (2) MANDT present  "
            "(3) AuthContext filters  (4) No Cartesian product  (5) LIMIT/max_rows guard  "
            "(6) JOIN keys exist  (7) Date range reasonable",
            font_name="Calibri", font_size=10, italic=True,
            color=RGBColor(0x66, 0x66, 0x66))

add_footer(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SELF-CRITIQUE / SELF-HEAL / SELF-IMPROVE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "06  Self-Critique / Self-Heal / Self-Improve",
                        "Autonomous SQL repair and continuous learning")

# Left — Self-Heal
add_rect(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(4.8),
         fill_color=RGBColor(0xF8, 0xF8, 0xFF),
         line_color=NRL_BLUE, line_width=Pt(1))
add_textbox(slide, Inches(0.55), Inches(1.45), Inches(5.8), Inches(0.4),
            "Phase 6a — Self-Healer (10 Error Codes, 6 Strategies)",
            font_name="Montserrat", font_size=14, bold=True, color=NRL_BLUE)
add_rect(slide, Inches(0.55), Inches(1.82), Inches(5.7), Inches(0.04), fill_color=NRL_YELLOW)

error_codes = [
    ("MANDT_MISSING",       "Add MANDT = '<client>' WHERE clause"),
    ("CARTESIAN_PRODUCT",   "Remove cross-join, add proper WHERE"),
    ("DIVISION_BY_ZERO",    "Wrap denominator in NULLIF(…, 0)"),
    ("TABLE_NOT_FOUND",     "Replace with DDIC-compliant alias"),
    ("INVALID_COLUMN",      "Qualify with table prefix"),
    ("SYNTAX_ERROR",        "Fix parentheses / missing commas"),
    ("SUBQUERY_JOIN_ERROR", "Align correlated subquery keys"),
    ("SAP_AUTH_BLOCK",      "Escalate to AuthContext gate"),
    ("EMPTY_RESULT",        "Relax filters, add NVL defaults"),
    ("ADD_NVL",             "Wrap nullable aggregations"),
]
for i, (code, fix) in enumerate(error_codes):
    y = Inches(1.92) + i * Inches(0.42)
    c = NRL_RED if i % 2 == 0 else RGBColor(0xDD, 0x33, 0x00)
    add_textbox(slide, Inches(0.6), y, Inches(1.8), Inches(0.38),
                code, font_name="Calibri", font_size=9, bold=True, color=c)
    add_textbox(slide, Inches(2.45), y, Inches(3.9), Inches(0.38),
                fix, font_name="Calibri", font_size=9, color=NRL_BLACK)

# Right — Self-Improver
add_rect(slide, Inches(6.8), Inches(1.35), Inches(6.1), Inches(2.5),
         fill_color=RGBColor(0xF8, 0xFF, 0xF8),
         line_color=RGBColor(0x00, 0x70, 0x00), line_width=Pt(1))
add_textbox(slide, Inches(6.95), Inches(1.45), Inches(5.8), Inches(0.4),
            "Phase 9 — Self-Improver (PSI Metric)",
            font_name="Montserrat", font_size=14, bold=True,
            color=RGBColor(0x00, 0x70, 0x00))
add_rect(slide, Inches(6.95), Inches(1.82), Inches(5.7), Inches(0.04), fill_color=NRL_YELLOW)

imp_rows = [
    ["Promote",   "≥5 consecutive successes",    "Pattern Stability Index (PSI) ≥ 0.85"],
    ["Demote",    "≥3 failures, ratio < 0.4",      "PSI drops below 0.4 threshold"],
    ["Ghost-Inject","Ad-hoc seen ≥3×, good score","Auto-promotes to named pattern"],
    ["Gotcha Log", "Heal-dependent patterns",      "Logged as edge-case rules"],
]
add_2col_table(slide, Inches(6.95), Inches(1.92), Inches(5.7), Inches(1.8),
               ["Action", "Trigger", "Effect"],
               [[r[0], r[1], r[2]] for r in imp_rows],
               col_widths=[Inches(1.1), Inches(2.3), Inches(2.3)])

# PSI Formula box
add_rect(slide, Inches(6.8), Inches(4.0), Inches(6.1), Inches(1.7),
         fill_color=NRL_DNAVY)
add_textbox(slide, Inches(6.95), Inches(4.1), Inches(5.8), Inches(0.35),
            "Pattern Stability Index (PSI) Formula",
            font_name="Montserrat", font_size=13, bold=True, color=NRL_YELLOW)
psi_items = [
    "PSI = (successes − failures) / (successes + failures + 1)",
    "Range: −1.0 (all failed) → +1.0 (all succeeded)",
    "Decisions: Promote ≥ 0.85 | Demote < 0.40 | Ghost ≥ 3 hits",
]
for i, item in enumerate(psi_items):
    add_textbox(slide, Inches(7.0), Inches(4.5) + i * Inches(0.35), Inches(5.7), Inches(0.32),
                f"\u2022  {item}", font_name="Calibri", font_size=12,
                color=NRL_WHITE)

add_footer(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SECURITY, AUTHCONTEXT & 7-LAYER SANDBOX
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "07  Security, AuthContext & 7-Layer Sandbox",
                        "Default-deny, layer-gated enforcement")

# Left — AuthContext
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.35),
            "SAPAuthContext — Role Scoping",
            font_name="Montserrat", font_size=15, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(0.4), Inches(1.63), Inches(5.8), Inches(0.04), fill_color=NRL_YELLOW)

auth_headers = ["Role", "Company Codes", "Denied Tables", "Masked Fields"]
auth_rows = [
    ["AP_CLERK",              "1000 / 1010",      "PA0008, MBEW",    "LFBK-BANKN, LFA1-STCD1"],
    ["PROCUREMENT_MANAGER_EU","2000 / 2010",      "PA0008",           "LFBK-BANKN"],
    ["CFO_GLOBAL",            "ALL",               "PA0008, PA0014",  "(none — FI access)"],
    ["HR_ADMIN",              "ALL",               "EKKO, EKPO, BSEG","(none — HR access)"],
]
add_2col_table(slide, Inches(0.4), Inches(1.73), Inches(5.8), Inches(2.0),
               auth_headers, auth_rows,
               col_widths=[Inches(1.6), Inches(1.2), Inches(1.6), Inches(1.6)],
               header_fill=NRL_DNAVY)

auth_constraints = [
    "get_where_clauses() — per-table WHERE injection",
    "MANDT presence — HARD BLOCK at Layer 3 (not soft warning)",
    "Role escalation mid-session — FORBIDDEN",
    "New role requires new session",
]
add_bullet_block(slide, Inches(0.4), Inches(3.85), Inches(5.8), Inches(1.4),
                 auth_constraints, font_size=12)

# Right — 7-Layer Sandbox
add_textbox(slide, Inches(6.7), Inches(1.3), Inches(6.2), Inches(0.35),
            "7-Layer Security Sandbox",
            font_name="Montserrat", font_size=15, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(6.7), Inches(1.63), Inches(6.2), Inches(0.04), fill_color=NRL_YELLOW)

sandbox = [
    ("7", "Network Isolation",     "K8s NetworkPolicy / docker --network=none",       "Env (TCP reject)",     NRL_DNAVY),
    ("6", "Threat Sentinel",        "6 behavioral engines — ENFORCING mode",          "Dynamic tighten",       RGBColor(0xCC, 0x00, 0x00)),
    ("5", "Result Masking",         "Post-execution column redaction",                 "cols → ***RESTRICTED***", RGBColor(0xCC, 0x55, 0x00)),
    ("4", "Table AuthContext",      "Hard block on denied_tables",                      "AuthorizationError",    RGBColor(0x80, 0x60, 0x00)),
    ("3", "SQL Permission Guard",   "Read-only + MANDT check — HARD BLOCK",            "PermissionError",       NRL_BLUE),
    ("2", "Dry-Run Harness",        "SELECT COUNT(*) trial before live",                "Self-Healer invoked",   RGBColor(0x00, 0x60, 0x80)),
    ("1", "Proxy Pattern",          "Credential isolation — no direct HANA contact",   "Request rejected",      RGBColor(0x00, 0x70, 0x00)),
]

for i, (num, name, desc, action, color) in enumerate(sandbox):
    y = Inches(1.73) + i * Inches(0.67)
    add_rect(slide, Inches(6.7), y, Inches(6.2), Inches(0.6),
             fill_color=RGBColor(0xF5, 0xF5, 0xFA),
             line_color=color, line_width=Pt(1.2))
    # Layer badge
    add_rect(slide, Inches(6.72), y + Inches(0.08), Inches(0.4), Inches(0.44),
             fill_color=color)
    add_textbox(slide, Inches(6.72), y + Inches(0.14), Inches(0.4), Inches(0.35),
                num, font_name="Montserrat", font_size=13, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.2), y + Inches(0.04), Inches(2.2), Inches(0.3),
                name, font_name="Montserrat", font_size=11, bold=True, color=color)
    add_textbox(slide, Inches(7.2), y + Inches(0.3), Inches(3.5), Inches(0.28),
                desc, font_name="Calibri", font_size=9,
                color=RGBColor(0x44, 0x44, 0x44))
    add_textbox(slide, Inches(10.7), y + Inches(0.14), Inches(2.1), Inches(0.3),
                f"\u2192 {action}", font_name="Calibri", font_size=9, italic=True,
                color=color)

add_footer(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THREAT SENTINEL & MESSAGE BUS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "07b  Threat Sentinel & Message Bus",
                        "Behavioral anomaly detection and inter-agent negotiation")

# Left — Sentinel
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.35),
            "Phase 7 — 6 Sentinel Behavioral Engines",
            font_name="Montserrat", font_size=14, bold=True, color=RGBColor(0xCC, 0x00, 0x00))
add_rect(slide, Inches(0.4), Inches(1.63), Inches(5.9), Inches(0.04), fill_color=NRL_RED)

sentinel_rows = [
    ["CROSS_MODULE_ESCALATION",  "Multi-hop graph traversal to out-of-scope tables"],
    ["SCHEMA_ENUMERATION",       "Bulk table discovery probes (>5 new tables/query)"],
    ["DENIED_TABLE_PROBE",       "Repeated attempts to access blocked tables"],
    ["DATA_EXFILTRATION",        "Unusually large result sets (>5,000 rows)"],
    ["TEMPORAL_INFERENCE",       "Restricted historical period access by HR/AP"],
    ["ROLE_IMPERSONATION",       "Sudden cross-domain shift mid-session"],
]
add_2col_table(slide, Inches(0.4), Inches(1.73), Inches(5.9), Inches(2.6),
               ["Engine", "Detects"],
               sentinel_rows,
               col_widths=[Inches(2.3), Inches(3.6)],
               header_fill=RGBColor(0xCC, 0x00, 0x00))

add_textbox(slide, Inches(0.4), Inches(4.45), Inches(5.9), Inches(0.3),
            "Sentinel Modes:  DISABLED  |  AUDIT  |  ENFORCING",
            font_name="Montserrat", font_size=12, bold=True, color=NRL_DNAVY)
mode_items = [
    "ENFORCING: dynamically adds denied_tables, expands masked_fields",
    "Dynamic Tightening: session tightness 0→3 escalates automatically",
    "6 message types: QUERY, RESPONSE, ASSERTION, CHALLENGE, NEGOTIATE, COMMIT",
]
add_bullet_block(slide, Inches(0.4), Inches(4.8), Inches(5.9), Inches(1.3), mode_items,
                 font_size=11)

# Right — Negotiation
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.35),
            "Phase 13b — Negotiation Protocol",
            font_name="Montserrat", font_size=14, bold=True, color=NRL_BLUE)
add_rect(slide, Inches(6.8), Inches(1.63), Inches(6.1), Inches(0.04), fill_color=NRL_YELLOW)

neg_states = ["ASSERTING", "CHALLENGING", "NEGOTIATING", "COMMITTED"]
state_cols = [NRL_BLUE, RGBColor(0xCC, 0x55, 0x00), NRL_YELLOW, RGBColor(0x00, 0x80, 0x00)]
for i, (st, col) in enumerate(zip(neg_states, state_cols)):
    x = Inches(6.9) + i * Inches(1.5)
    add_rect(slide, x, Inches(1.73), Inches(1.35), Inches(0.4), fill_color=col)
    add_textbox(slide, x, Inches(1.77), Inches(1.35), Inches(0.35),
                st, font_name="Montserrat", font_size=9, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + Inches(1.35), Inches(1.77), Inches(0.15), Inches(0.35),
                    "\u2192", font_name="Calibri", font_size=12,
                    color=NRL_BLACK, align=PP_ALIGN.CENTER)

res_headers = ["Strategy", "Use Case", "Notes"]
res_rows = [
    ["DOMAIN_SPECIFIC_OVERRIDE", "Numeric conflicts on specialist data", "DEFAULT"],
    ["HIGHEST_AUTHORITY",         "Planner weight 8 outranks all others", "Planner owns"],
    ["AVERAGE",                   "Balanced numeric merge across agents", "Cross-domain"],
    ["MAJORITY_VOTE",             "Categorical value agreement",         "Value-based"],
    ["BATNA_FALLBACK",            "No resolution — use best alternative", "Last resort"],
    ["PLANNER_ARBITRATION",       "Planner resolves unresolvable conflicts", "Final fallback"],
]
add_2col_table(slide, Inches(6.8), Inches(2.23), Inches(6.1), Inches(2.4),
               res_headers, res_rows,
               col_widths=[Inches(2.6), Inches(2.0), Inches(1.5)],
               header_fill=NRL_DNAVY)

add_footer(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GRAPH RAG & MEMGRAPH
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "08  Graph RAG & Memgraph",
                        "Dual-mode graph engine with 14 meta-path templates")

# Stats bar
stats_bar = [
    ("114", "DDIC Tables"),
    ("137", "Directed Edges"),
    ("151", "With Birectional"),
    ("97",  "Cross-Module Bridges"),
    ("14",  "Meta-Path Templates"),
    ("2",   "Execution Modes"),
]
bar_y = Inches(1.35)
add_rect(slide, Inches(0.4), bar_y, Inches(12.5), Inches(0.9),
         fill_color=NRL_DNAVY)
for i, (val, lbl) in enumerate(stats_bar):
    x = Inches(0.5) + i * Inches(2.1)
    add_textbox(slide, x, bar_y + Inches(0.05), Inches(1.8), Inches(0.5),
                val, font_name="Montserrat", font_size=22, bold=True,
                color=NRL_YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, bar_y + Inches(0.52), Inches(1.8), Inches(0.3),
                lbl, font_name="Calibri", font_size=10,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)

# Left — Techniques
add_textbox(slide, Inches(0.4), Inches(2.4), Inches(6.0), Inches(0.35),
            "Graph RAG Techniques",
            font_name="Montserrat", font_size=14, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(0.4), Inches(2.72), Inches(6.0), Inches(0.04), fill_color=NRL_YELLOW)

techs = [
    "BFS Shortest-Path — efficient single-hop discovery",
    "All-Simple-Paths (max depth 5) — scored for 1:N / bridge bonuses",
    "Meta-Path Templates — 14 named JOIN chains across modules",
    "Steiner Tree — minimum connection across ≥3 entity tables",
    "TemporalGraphRAG — fiscal-year / key-date / range filters",
    "Node2Vec 64-dim embeddings — hub / bridge / authority / spoke roles",
    "Dual-mode: Memgraph (authoritative) + NetworkX mirror (fallback)",
]
add_bullet_block(slide, Inches(0.4), Inches(2.82), Inches(6.0), Inches(2.5), techs,
                 font_size=12)

# Right — Meta-Path examples
add_textbox(slide, Inches(6.8), Inches(2.4), Inches(6.1), Inches(0.35),
            "14 Meta-Path Templates (examples)",
            font_name="Montserrat", font_size=14, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(6.8), Inches(2.72), Inches(6.1), Inches(0.04), fill_color=NRL_YELLOW)

meta_paths = [
    ("procure_to_pay",    "LFA1 \u2192 EKKO \u2192 EKPO \u2192 MSEG \u2192 BKPF \u2192 BSEG"),
    ("order_to_cash",     "KNA1 \u2192 VBAK \u2192 LIKP \u2192 VBRK \u2192 BSEG"),
    ("vendor_master",     "LFA1 \u2192 LFB1 \u2192 LFBK \u2192 ADRC"),
    ("material_stock",    "MARA \u2192 MARD \u2192 MBEW + MSKA/MSLB"),
    ("inspection_lot",    "QALS \u2192 QAVE \u2192 QAMV \u2192 MAPL"),
    ("project_costs",     "PRPS \u2192 COSP/COSS \u2192 AFKO/AFVC"),
    ("transportation",    "LIKP \u2192 VTTK \u2192 VTTS \u2192 TVRO"),
]
for i, (name, path) in enumerate(meta_paths):
    y = Inches(2.82) + i * Inches(0.52)
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i % 2 == 0 else NRL_WHITE
    add_rect(slide, Inches(6.8), y, Inches(6.1), Inches(0.47),
             fill_color=bg, line_color=NRL_BLUE, line_width=Pt(0.5))
    add_textbox(slide, Inches(6.9), y + Inches(0.07), Inches(1.8), Inches(0.32),
                name, font_name="Calibri", font_size=10, bold=True, color=NRL_BLUE)
    add_textbox(slide, Inches(8.7), y + Inches(0.07), Inches(4.0), Inches(0.32),
                path, font_name="Calibri", font_size=9, color=NRL_BLACK)

# Dual mode note
add_rect(slide, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.55),
         fill_color=RGBColor(0xE8, 0xF0, 0xFF),
         line_color=NRL_BLUE, line_width=Pt(0.75))
add_textbox(slide, Inches(0.55), Inches(6.08), Inches(12.2), Inches(0.42),
            "Dual-Mode: Memgraph authoritative. If unreachable \u2192 NetworkX mirror transparently. "
            "init_schema.cql is single source of truth for both. Startup warning logged as [STARTUP] if Memgraph is down.",
            font_name="Calibri", font_size=11, color=NRL_BLACK)

add_footer(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PERSISTENT MEMORY & DIALOG MANAGER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "09  Persistent Memory & Dialog Manager",
                        "6 stores + session management")

# Memory stores
stores = [
    ("query_history",      "query_history.jsonl",     "Per-query audit trail — role, SQL, latency, verdict"),
    ("pattern_success",    "pattern_success.json",     "Promoted patterns (PSI \u2265 0.85, \u22655 consecutive)"),
    ("pattern_failures",   "pattern_failures.json",    "Demoted patterns (PSI < 0.4 or \u22653 failures)"),
    ("schema_discoveries", "schema_discoveries.json",  "DDIC fallback hits — table/field cache"),
    ("gotchas",            "gotchas.json",              "Edge-case ledger — heal-dependent rules"),
    ("user_preferences",   "user_preferences.json",    "Per-role format, max_rows, language pref"),
]
for i, (store, file, desc) in enumerate(stores):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.5)
    y = Inches(1.4) + row * Inches(1.5)
    add_rect(slide, x, y, Inches(6.2), Inches(1.3),
             fill_color=RGBColor(0xF5, 0xF8, 0xFF),
             line_color=NRL_BLUE, line_width=Pt(1))
    add_rect(slide, x, y, Inches(0.25), Inches(1.3), fill_color=NRL_BLUE)
    add_textbox(slide, x + Inches(0.35), y + Inches(0.1), Inches(5.7), Inches(0.35),
                store, font_name="Montserrat", font_size=13, bold=True, color=NRL_BLUE)
    add_textbox(slide, x + Inches(0.35), y + Inches(0.42), Inches(5.7), Inches(0.28),
                file, font_name="Calibri", font_size=10, italic=True,
                color=RGBColor(0x88, 0x00, 0x88))
    add_textbox(slide, x + Inches(0.35), y + Inches(0.72), Inches(5.7), Inches(0.45),
                desc, font_name="Calibri", font_size=11, color=NRL_BLACK)

# Storage location note
add_rect(slide, Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.5),
         fill_color=NRL_DNAVY)
add_textbox(slide, Inches(0.55), Inches(5.98), Inches(12.2), Inches(0.38),
            "Location: ~/.openclaw/workspace/memory/sap_sessions/{tenant_id}/{role}/  "
            "All writes tenant-scoped. Cross-tenant reads fail closed.",
            font_name="Calibri", font_size=12, bold=True,
            color=NRL_YELLOW)

# Dialog Manager
add_textbox(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.35),
            "Dialog Manager — 5 Clarification Types:  SCOPE  |  ENTITY  |  TIME_RANGE  |  METRIC  |  DOMAIN",
            font_name="Montserrat", font_size=13, bold=True, color=NRL_DNAVY)

add_footer(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — NFRs & COMPLIANCE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "10  Non-Functional Requirements & Compliance",
                        "Performance targets, scalability, availability, Willison scorecard")

# Left — NFR table
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.35),
            "Performance & Scalability Targets",
            font_name="Montserrat", font_size=14, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(0.4), Inches(1.63), Inches(5.9), Inches(0.04), fill_color=NRL_YELLOW)

nfr_headers = ["Metric", "Target"]
nfr_rows = [
    ["p50 E2E latency (cached)",        "\u2264 180 ms"],
    ["p95 E2E latency (full pipeline)",  "\u2264 300 ms @ concurrency 10"],
    ["Memgraph path query (depth \u22643)", "\u2264 40 ms"],
    ["Qdrant top-K (K=5)",               "\u2264 25 ms"],
    ["Self-heal retry budget",           "\u2264 2 retries / query"],
    ["Target availability",              "99.5% monthly uptime"],
    ["Celery workers (KEDA HPA)",        "min=2, max=20, cooldown=120s"],
    ["Trajectory log retention",         "\u2265 90 days"],
]
add_2col_table(slide, Inches(0.4), Inches(1.73), Inches(5.9), Inches(3.0),
               nfr_headers, nfr_rows,
               col_widths=[Inches(3.5), Inches(2.4)],
               header_fill=NRL_DNAVY)

# Right — Willison Scorecard
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.35),
            "Simon Willison 14-Pattern Scorecard",
            font_name="Montserrat", font_size=14, bold=True,
            color=RGBColor(0x00, 0x60, 0x00))
add_rect(slide, Inches(6.8), Inches(1.63), Inches(6.1), Inches(0.04), fill_color=NRL_YELLOW)

# Score badge
add_rect(slide, Inches(6.8), Inches(1.73), Inches(6.1), Inches(0.65),
         fill_color=RGBColor(0xE8, 0xFF, 0xE8),
         line_color=RGBColor(0x00, 0x70, 0x00), line_width=Pt(1))
add_textbox(slide, Inches(6.95), Inches(1.8), Inches(1.5), Inches(0.55),
            "11 / 14", font_name="Montserrat", font_size=24, bold=True,
            color=RGBColor(0x00, 0x70, 0x00), align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.5), Inches(1.88), Inches(4.2), Inches(0.4),
            "11 \u2705 PASS  |  3 \u26a0 NEEDS ATTENTION",
            font_name="Calibri", font_size=12, bold=True, color=RGBColor(0x00, 0x70, 0x00))

willison = [
    ("1  Validate untrusted input",           True),
    ("2  Allow-lists over deny-lists",         True),
    ("3  Sandbox dangerous operations",        True),
    ("4  Least privilege",                     True),
    ("5  Defense in depth",                    True),
    ("6  Red/Green TDD for prompts",           False),
    ("7  Observability",                       True),
    ("8  Rate limiting",                       False),
    ("9  Secrets out of prompts",              True),
    ("10 Audit logs",                          True),
    ("11 Deterministic replays (Traj. Log)",   True),
    ("12 Human-in-the-loop escalation",        False),
    ("13 Capability-based tool access",        True),
    ("14 Memory isolation",                    True),
]
for i, (pattern, passed) in enumerate(willison):
    col = i % 2
    row = i // 2
    x = Inches(6.8) + col * Inches(3.05)
    y = Inches(2.48) + row * Inches(0.47)
    bg = RGBColor(0xE8, 0xFF, 0xE8) if passed else RGBColor(0xFF, 0xF0, 0xE0)
    line = RGBColor(0x00, 0x70, 0x00) if passed else NRL_RED
    add_rect(slide, x, y, Inches(2.95), Inches(0.42),
             fill_color=bg, line_color=line, line_width=Pt(0.5))
    icon = "\u2705" if passed else "\u26a0"
    add_textbox(slide, x + Inches(0.05), y + Inches(0.07), Inches(2.85), Inches(0.3),
                f"{icon}  {pattern}", font_name="Calibri", font_size=10,
                color=RGBColor(0x22, 0x22, 0x22))

add_footer(slide)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PENDING REQUIREMENTS & ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "11  Pending Requirements & Roadmap",
                        "P0–P3 priorities for v1.0 completion")

# Priority bands
priority_bands = [
    ("P0 — Critical",  NRL_RED,   [
        ("P0-HANA", "Switch Phase 6 executor from mock to real hdbcli pool", "2026-05"),
    ]),
    ("P1 — High",     RGBColor(0xCC, 0x77, 0x00), [
        ("P1-BAPI",    "Write-back via BAPI for LFA1, MARA",                  "2026-Q3"),
        ("M7-POOL",    "hdbcli connection-pool sizing (currently 10)",       "2026-05"),
        ("W-CI-GATE",  "Willison Pattern 6 golden set as merge gate",          "2026-05"),
        ("R-RATELIMIT","Per-tenant rate limiting (Willison Pattern 8)",         "2026-05"),
    ]),
    ("P2 — Medium",   NRL_BLUE, [
        ("P2-BENCH",   "Wire eval benchmark into nightly CI",                 "2026-05"),
        ("W-HITL",     "Human-in-the-loop escalation UI (Willison P12)",       "2026-Q3"),
    ]),
    ("P3 — Low",      RGBColor(0x00, 0x70, 0x00), [
        ("P3-LOAD",    "500 QPS load test sign-off with KEDA autoscale",      "2026-06"),
    ]),
]

for band_i, (band_name, band_color, items) in enumerate(priority_bands):
    y = Inches(1.35) + band_i * Inches(1.45)
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(1.3),
             fill_color=RGBColor(0xF8, 0xF8, 0xFF),
             line_color=band_color, line_width=Pt(1.5))
    add_rect(slide, Inches(0.4), y, Inches(1.6), Inches(1.3), fill_color=band_color)
    add_textbox(slide, Inches(0.45), y + Inches(0.45), Inches(1.5), Inches(0.4),
                band_name, font_name="Montserrat", font_size=12, bold=True,
                color=NRL_WHITE, align=PP_ALIGN.CENTER)
    for j, (tid, desc, date) in enumerate(items):
        ix = Inches(2.1) + j * Inches(2.6)
        add_textbox(slide, ix, y + Inches(0.1), Inches(2.4), Inches(0.32),
                    tid, font_name="Montserrat", font_size=11, bold=True,
                    color=band_color)
        add_textbox(slide, ix, y + Inches(0.42), Inches(2.4), Inches(0.42),
                    desc, font_name="Calibri", font_size=10,
                    color=NRL_BLACK)
        add_textbox(slide, ix, y + Inches(0.88), Inches(2.4), Inches(0.28),
                    f"\u231a {date}", font_name="Calibri", font_size=10,
                    italic=True, color=RGBColor(0x88, 0x88, 0x88))

add_footer(slide)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — USER CLASSES & API
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_WHITE)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=NRL_YELLOW)
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=NRL_BLUE)
add_title_slide_header(slide, "User Classes & API Interface",
                        "Built-in roles and REST API contract")

# Left — User Classes
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.35),
            "4 Built-in SAP Roles (extensible via SAPAuthContext)",
            font_name="Montserrat", font_size=13, bold=True, color=NRL_DNAVY)
add_rect(slide, Inches(0.4), Inches(1.63), Inches(5.9), Inches(0.04), fill_color=NRL_YELLOW)
roles = [
    ("AP_CLERK",              "US (1000/1010)",    "PA0008, MBEW",          "LFBK-BANKN, LFA1-STCD1"),
    ("PROCUREMENT_MANAGER_EU","EU (2000/2010)",    "PA0008",                "LFBK-BANKN"),
    ("CFO_GLOBAL",            "ALL companies",      "PA0008, PA0014",        "(none)"),
    ("HR_ADMIN",              "ALL companies",      "EKKO, EKPO, BSEG",     "(none)"),
]
add_2col_table(slide, Inches(0.4), Inches(1.73), Inches(5.9), Inches(2.0),
               ["Role", "Company Codes", "Denied Tables", "Masked Fields"],
               roles,
               col_widths=[Inches(1.8), Inches(1.2), Inches(1.5), Inches(1.6)],
               header_fill=NRL_DNAVY)

# Constraints
add_textbox(slide, Inches(0.4), Inches(3.9), Inches(5.9), Inches(0.3),
            "Hard Constraints (v1.0)",
            font_name="Montserrat", font_size=13, bold=True, color=NRL_DNAVY)
constraints = [
    "C1 — Read-only. No INSERT / UPDATE / DELETE / DDL to HANA.",
    "C2 — MANDT mandatory. Every SQL includes MANDT = '<client>'.",
    "C3 — Network isolation. Workers have no outbound internet.",
    "C4 — No credential in prompt. LLM receives no HANA connection strings.",
    "C5 — Default-deny security. Unknown tables/columns are blocked.",
    "C6 — Multi-tenant isolation. All state scoped by TENANT_ID.",
]
add_bullet_block(slide, Inches(0.4), Inches(4.25), Inches(5.9), Inches(2.4), constraints,
                 font_size=11)

# Right — API
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.35),
            "REST API — POST /api/v1/chat/master-data",
            font_name="Montserrat", font_size=13, bold=True, color=NRL_BLUE)
add_rect(slide, Inches(6.8), Inches(1.63), Inches(6.1), Inches(0.04), fill_color=NRL_YELLOW)

add_rect(slide, Inches(6.8), Inches(1.73), Inches(6.1), Inches(2.1),
         fill_color=RGBColor(0xF0, 0xF4, 0xFF))
add_textbox(slide, Inches(6.9), Inches(1.8), Inches(5.9), Inches(0.3),
            "Request Body:", font_name="Montserrat", font_size=11, bold=True,
            color=NRL_BLUE)
req_fields = [
    "tenant_id: string",
    "role: AP_CLERK | PROCUREMENT_MANAGER_EU | CFO_GLOBAL | HR_ADMIN",
    "session_id: uuid",
    "query: string",
    "domain_hint: business_partner | material | purchasing | ... | null",
    "verbose: boolean (default false)",
]
add_bullet_block(slide, Inches(6.9), Inches(2.1), Inches(5.9), Inches(1.6), req_fields,
                 font_size=10)

add_rect(slide, Inches(6.8), Inches(3.95), Inches(6.1), Inches(2.4),
         fill_color=RGBColor(0xF5, 0xFF, 0xF5))
add_textbox(slide, Inches(6.9), Inches(4.02), Inches(5.9), Inches(0.3),
            "Response Fields:", font_name="Montserrat", font_size=11, bold=True,
            color=RGBColor(0x00, 0x70, 0x00))
resp_fields = [
    "answer: string (natural language)",
    "sql_executed: string",
    "tables_used: string[]",
    "critique_score: float (0.0–1.0)",
    "confidence_score: object (composite + 6 signals)",
    "rows: integer",
    "columns_masked: string[]",
    "trajectory_id: uuid",
    "latency_ms: integer",
    "sentinel: object (verdict + engines triggered)",
]
add_bullet_block(slide, Inches(6.9), Inches(4.35), Inches(5.9), Inches(2.0), resp_fields,
                 font_size=10)

add_footer(slide)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU / CLOSING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NRL_DNAVY)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=NRL_YELLOW)

# Center
add_textbox(slide, Inches(1), Inches(2.0), SLIDE_W - Inches(2), Inches(1.2),
            "Thank You",
            font_name="Montserrat", font_size=56, bold=True,
            color=NRL_WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5), Inches(3.3), Inches(3.333), Inches(0.05), fill_color=NRL_YELLOW)

add_textbox(slide, Inches(1), Inches(3.5), SLIDE_W - Inches(2), Inches(0.6),
            "Know Your SAP Masters — Agentic RAG System",
            font_name="Calibri", font_size=22,
            color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.15), SLIDE_W - Inches(2), Inches(0.5),
            "Functional Requirement Specification v1.0  |  KYSM-FRS-001",
            font_name="Calibri", font_size=14,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.65), SLIDE_W - Inches(2), Inches(0.5),
            "Numaligarh Refinery Limited  |  KYSM Platform Team",
            font_name="Calibri", font_size=14,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.15), SLIDE_W - Inches(2), Inches(0.5),
            "2026-04-16",
            font_name="Calibri", font_size=14,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# Bottom tagline
add_rect(slide, 0, SLIDE_H - Inches(0.7), SLIDE_W, Inches(0.7), fill_color=RGBColor(0x00, 0x1A, 0x4A))
add_textbox(slide, Inches(1), SLIDE_H - Inches(0.62), SLIDE_W - Inches(2), Inches(0.5),
            "Energy in Motion",
            font_name="Montserrat", font_size=18, bold=True,
            color=NRL_YELLOW, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
output_path = "C:/Users/vishnu/.openclaw/workspace/NRL_KYSM_FRS_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
