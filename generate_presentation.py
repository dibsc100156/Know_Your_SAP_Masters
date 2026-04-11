"""
Know Your SAP Masters — 5-Pillar RAG Architecture
High-Fidelity PPTX Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree

# ── Brand Palette ──────────────────────────────────────────────────────────────
C_DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy – backgrounds
C_NAVY        = RGBColor(0x1A, 0x37, 0x6E)   # mid navy – headers
C_ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xC2)   # SAP-like accent
C_ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)   # highlight / callout
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF5)   # light fill
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF7)
C_MID_GRAY    = RGBColor(0x8A, 0x9B, 0xA8)
C_TEXT_DARK   = RGBColor(0x1C, 0x2D, 0x4A)
C_GREEN       = RGBColor(0x00, 0xB0, 0x6A)   # Pillar 1
C_PURPLE      = RGBColor(0x6B, 0x47, 0xD0)   # Pillar 2
C_ORANGE      = RGBColor(0xE8, 0x7D, 0x1A)   # Pillar 3
C_RED         = RGBColor(0xCC, 0x29, 0x2A)   # Pillar 4
C_TEAL        = RGBColor(0x00, 0x96, 0x88)   # Pillar 5
C_ACCENT_RED  = RGBColor(0xCC, 0x29, 0x2A)   # accent red

PILLAR_COLORS = [C_GREEN, C_PURPLE, C_ORANGE, C_RED, C_TEAL]
PILLAR_NAMES  = [
    "Pillar 1 – Role-Aware RAG",
    "Pillar 2 – Agentic RAG",
    "Pillar 3 – Schema RAG",
    "Pillar 4 – SQL RAG",
    "Pillar 5 – Graph RAG",
]

# ── Slide Dimensions (Widescreen 16:9) ─────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ═══════════════════════════════════════════════════════════════════════════════
# HELPER UTILITIES
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    elif not line_rgb:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, italic=False,
                color=C_TEXT_DARK, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox


def add_para(tf, text, font_size=12, bold=False, italic=False,
             color=C_TEXT_DARK, align=PP_ALIGN.LEFT, space_before=0,
             font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return p


def header_bar(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill_rgb=C_DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), W, Pt(4), fill_rgb=C_ACCENT_BLUE)
    t = add_textbox(slide, title,
                    Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.65),
                    font_size=24, bold=True, color=C_WHITE,
                    align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.35),
                    font_size=13, bold=False, color=C_ACCENT_BLUE,
                    align=PP_ALIGN.LEFT)


def footer_bar(slide, text="Know Your SAP Masters — 5-Pillar RAG Architecture  |  Confidential"):
    add_rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill_rgb=C_DARK_NAVY)
    add_textbox(slide, text,
                Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
                font_size=8, color=C_MID_GRAY, align=PP_ALIGN.LEFT)


def section_divider(slide, number, title, subtitle=""):
    """Full-bleed dark divider slide."""
    add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
    # Large number watermark
    add_textbox(slide, f"0{number}",
                Inches(0.5), Inches(1.2), Inches(4), Inches(3.5),
                font_size=160, bold=True,
                color=RGBColor(0x1A, 0x37, 0x6E),
                align=PP_ALIGN.LEFT)
    # Accent line
    add_rect(slide, Inches(4.5), Inches(2.9), Inches(0.06), Inches(1.8),
             fill_rgb=C_ACCENT_BLUE)
    add_textbox(slide, title,
                Inches(4.8), Inches(2.7), Inches(8), Inches(1.2),
                font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(4.8), Inches(4.0), Inches(8), Inches(0.6),
                    font_size=16, bold=False, color=C_ACCENT_BLUE,
                    align=PP_ALIGN.LEFT)
    footer_bar(slide)


def pillar_icon(slide, number, left, top, size=Inches(0.5)):
    """Colored circle with pillar number."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PILLAR_COLORS[number - 1]
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(number)
    run.font.bold  = True
    run.font.size  = Pt(16)
    run.font.color.rgb = C_WHITE
    run.font.name  = "Calibri"


def add_bullet_box(slide, bullets, left, top, width, height,
                    title=None, title_color=C_ACCENT_BLUE,
                    bg_rgb=None, padding=Inches(0.15)):
    """Rounded-ish box with optional title and bullets."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    if bg_rgb:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_rgb
    else:
        box.fill.background()
    box.line.color.rgb = RGBColor(0xD0, 0xDA, 0xE4)
    box.line.width = Pt(0.75)

    txBox = slide.shapes.add_textbox(
        left + padding, top + padding,
        width - padding * 2, height - padding * 2)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for b in bullets:
        if first and title:
            p = tf.paragraphs[0]
            first = False
        elif first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(b, tuple):
            run.text = f"• {b[0]}"
            run.font.size = Pt(b[1] if len(b) > 1 else 11)
            if len(b) > 2 and b[2]:
                run.font.bold = True
                run.font.color.rgb = b[2]
            else:
                run.font.color.rgb = C_TEXT_DARK
        else:
            run.text = f"• {b}"
            run.font.size = Pt(11)
            run.font.color.rgb = C_TEXT_DARK
        run.font.name = "Calibri"
    return box


def table_style(table, header_rgb, row_rgb=None):
    """Apply banding to a pptx table."""
    for row_i, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            if row_i == 0:
                cell.fill.fore_color.rgb = header_rgb
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.bold  = True
                        run.font.color.rgb = C_WHITE
                        run.font.size = Pt(10)
            else:
                if row_i % 2 == 0:
                    cell.fill.fore_color.rgb = row_rgb or C_LIGHT_GRAY
                else:
                    cell.fill.fore_color.rgb = C_WHITE
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = C_TEXT_DARK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)

# Right accent panel
add_rect(slide, Inches(9.2), 0, Inches(4.13), H, fill_rgb=C_NAVY)
add_rect(slide, Inches(9.13), 0, Pt(6), H, fill_rgb=C_ACCENT_BLUE)

# "5-PILLAR RAG" label
lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.55), Inches(1.5), Inches(2.4), Inches(0.42))
lbl.fill.solid(); lbl.fill.fore_color.rgb = C_ACCENT_BLUE
lbl.line.fill.background()
tf = lbl.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "ENTERPRISE AI ARCHITECTURE"
r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=C_WHITE; r.font.name="Calibri"

# Main title
add_textbox(slide, "Know Your\nSAP Masters",
            Inches(0.5), Inches(2.1), Inches(8.4), Inches(2.2),
            font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide, "5-Pillar RAG Architecture for SAP S/4 HANA",
            Inches(0.5), Inches(4.35), Inches(8.4), Inches(0.55),
            font_size=20, bold=False, color=C_ACCENT_BLUE, align=PP_ALIGN.LEFT)

# Gold divider line
add_rect(slide, Inches(0.5), Inches(5.05), Inches(3.5), Pt(3),
         fill_rgb=C_ACCENT_GOLD)

# Tagline
add_textbox(slide, "From Vendor Chatbot → Enterprise Master Data Intelligence",
            Inches(0.5), Inches(5.2), Inches(8.4), Inches(0.45),
            font_size=13, italic=True, color=C_MID_GRAY, align=PP_ALIGN.LEFT)

# Right panel pillars
for i, (name, col) in enumerate(zip(PILLAR_NAMES, PILLAR_COLORS)):
    y = Inches(1.0) + i * Inches(1.2)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(9.55), y, Inches(3.3), Inches(0.85))
    pill.fill.solid(); pill.fill.fore_color.rgb = col
    pill.line.fill.background()
    tf2 = pill.text_frame; tf2.word_wrap = False
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run()
    r2.text = name.replace("Pillar ", "P").replace(" – ", "\n")
    r2.font.size = Pt(9); r2.font.bold = True
    r2.font.color.rgb = C_WHITE; r2.font.name = "Calibri"

footer_bar(slide, "Know Your SAP Masters — 5-Pillar RAG  |  Confidential")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "The Problem", "Why generic LLMs fail on SAP S/4 HANA")

problems = [
    ("90,000+ Tables",
     "SAP's DDIC is vast. Natural language queries cannot reliably map 'vendor' to LFA1 or 'pricing' to KONP without domain-specific retrieval."),
    ("SAP SQL Idioms",
     "LLMs hallucinate critical patterns: MANDT filters, deletion flags (LOEVM, LVORM), and index-aware join order (BUKRS first)."),
    ("Cross-Module Joins",
     "Enterprise questions span modules — 'Which vendors supply materials in Cost Center X?' requires MARA→EINA→LFA1→CSKS traversal."),
    ("Authorization Complexity",
     "SAP's object-based security (BUKRS, WERKS, EKORG) cannot be delegated to an LLM's internal reasoning. It must be a hard constraint."),
    ("No Audit Trail",
     "Every SQL-generating interaction must be logged with user identity, generated SQL, and auth context for SOX compliance."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.35) + i * Inches(1.18)
    add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.9),
             fill_rgb=C_ACCENT_RED)
    add_textbox(slide, title,
                Inches(0.65), y, Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=C_ACCENT_RED)
    add_textbox(slide, desc,
                Inches(0.65), y + Inches(0.32), Inches(11.8), Inches(0.58),
                font_size=10.5, color=C_TEXT_DARK)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – 5 PILLAR OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "The 5-Pillar Architecture", "A layered defense for enterprise-grade SAP SQL generation")

pillars = [
    (C_GREEN,  "Pillar 1", "Role-Aware RAG",
     ["Security perimeter & SOX compliance",
      "AuthContext enforcement at 3 layers",
      "Column-level masking & row filtering",
      "Immutable audit ledger"]),
    (C_PURPLE, "Pillar 2", "Agentic RAG",
     ["Autonomous orchestrator (ReAct loop)",
      "Multi-step planning & decomposition",
      "Self-correction on SQL errors",
      "Intent routing across modules"]),
    (C_ORANGE, "Pillar 3", "Schema RAG",
     ["Dynamic SAP DDIC metadata retrieval",
      "Qdrant vector search by module",
      "CDS View abstraction layer",
      "Covers 90,000+ tables"]),
    (C_RED,    "Pillar 4", "SQL RAG",
     ["Few-shot proven SAP SQL patterns",
      "MANDT / LOEVM idiom injection",
      "Performance-aware query library",
      "Auto-promotion of validated queries"]),
    (C_TEAL,   "Pillar 5", "Graph RAG",
     ["NetworkX cross-module traversal",
      "MARA→EINA→LFA1 join path discovery",
      "Shortest-path join sequencing",
      "Multi-hop Q&A (Phase 2)"]),
]

col_w = Inches(2.48)
for i, (col, num, name, bullets) in enumerate(pillars):
    x = Inches(0.3) + i * col_w
    # Card background
    add_rect(slide, x, Inches(1.3), col_w - Inches(0.1), Inches(5.9),
             fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0,0xDA,0xE4), line_width_pt=0.5)
    # Colored header strip
    add_rect(slide, x, Inches(1.3), col_w - Inches(0.1), Inches(0.7),
             fill_rgb=col)
    # Pillar number
    add_textbox(slide, num,
                x + Inches(0.1), Inches(1.35), col_w - Inches(0.3), Inches(0.3),
                font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # Pillar name
    add_textbox(slide, name,
                x + Inches(0.1), Inches(1.6), col_w - Inches(0.3), Inches(0.35),
                font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # Bullets
    for j, b in enumerate(bullets):
        y = Inches(2.15) + j * Inches(1.2)
        add_rect(slide, x + Inches(0.12), y, Pt(5), Pt(5),
                 fill_rgb=col)
        add_textbox(slide, b,
                    x + Inches(0.25), y - Inches(0.04), col_w - Inches(0.4), Inches(1.0),
                    font_size=9.5, color=C_TEXT_DARK)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – SECTION DIVIDER: PILLAR 1
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_divider(slide, 1, "Role-Aware RAG", "Pillar 1 — Security Mesh & SOX Compliance")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – ROLE-AWARE RAG DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Role-Aware RAG — 3-Layer Security Mesh",
           "Every query is intercepted, validated, and audited before execution")

# Left: layers
layers = [
    (C_GREEN, "Layer 1", "Prompt Injection",
     "AuthContext hardcoded into LLM system prompt.\n"
     "Defines BUKRS/WERKS allowed values and masked columns.\n"
     "LLM cannot 'forget' — it's structurally enforced."),
    (RGBColor(0x00,0x7A,0xC2), "Layer 2", "SQL Validation",
     "sqlglot parses generated SQL pre-execution.\n"
     "Verifies WHERE clauses exist for restricted fields.\n"
     "Blocks INSERT/UPDATE/DELETE and restricted columns."),
    (RGBColor(0xF5,0xA6,0x23), "Layer 3", "Response Masking",
     "Post-execution field redaction.\n"
     "Bank accounts, SSNs, salary data stripped.\n"
     "Final markdown sanitized before user display."),
]

for i, (col, layer, title, desc) in enumerate(layers):
    x = Inches(0.35) + i * Inches(4.25)
    add_rect(slide, x, Inches(1.35), Inches(4.0), Inches(1.2),
             fill_rgb=col)
    add_textbox(slide, f"{layer} — {title}",
                x + Inches(0.15), Inches(1.4), Inches(3.7), Inches(0.35),
                font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide, desc,
                x + Inches(0.15), Inches(1.78), Inches(3.7), Inches(0.72),
                font_size=9.5, color=C_WHITE, italic=True)

# Auth objects table
auth_data = [
    ["Domain", "SAP Auth Object", "Maps to Column"],
    ["Finance", "F_BKPF_BUK", "BUKRS"],
    ["Material", "M_MATE_WRK", "WERKS"],
    ["Purchasing", "M_BEST_EKO", "EKORG"],
    ["Sales", "V_VBAK_VKO", "VKORG"],
    ["Business Partner", "F_LFA1_BUK / F_KNA1_BUK", "BUKRS"],
]

tbl = slide.shapes.add_table(
    len(auth_data), 3,
    Inches(0.35), Inches(2.75), Inches(12.6), Inches(2.2)
).table
col_widths = [Inches(3.5), Inches(5.0), Inches(4.1)]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = cw
for ri, row_data in enumerate(auth_data):
    for ci, val in enumerate(row_data):
        cell = tbl.cell(ri, ci)
        cell.text = val
table_style(tbl, C_GREEN, C_LIGHT_GRAY)

add_textbox(slide, "SAP Authorization Object Mapping",
            Inches(0.35), Inches(2.62), Inches(12.6), Inches(0.25),
            font_size=10, bold=True, color=C_TEXT_DARK)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – SECTION DIVIDER: PILLAR 2
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_divider(slide, 2, "Agentic RAG", "Pillar 2 — Autonomous Reasoning & Orchestration")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – AGENTIC RAG DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Agentic RAG — ReAct Orchestration Loop",
           "The autonomous brain that plans, validates, self-corrects, and delivers")

# Flow diagram (manual boxes)
flow_steps = [
    (C_PURPLE, "USER\nQUERY",      "Intent Classification\n& Decomposition"),
    (RGBColor(0x4A,0x35,0x9E), "AUTH\nCHECK",   "get_auth_context()\nRole-Aware Pillar 1"),
    (C_ORANGE, "SCHEMA\nRAG",      "search_schema()\nPillar 3 — Qdrant"),
    (C_RED,    "SQL\nRAG",         "get_sql_patterns()\nPillar 4 — Few-shot"),
    (C_TEAL,   "GRAPH\nRAG",       "traverse_graph()\nPillar 5 — NetworkX"),
    (RGBColor(0x00,0x7A,0xC2), "SQL\nGEN",       "LLM → HANA SQL\n+ Auth Filters"),
    (C_GREEN,  "EXECUTE\n& VALIDATE","execute_hana_sql()\nSelf-Correction Loop"),
    (RGBColor(0xF5,0xA6,0x23), "FINAL\nANSWER",  "Natural Language\nResponse to User"),
]

box_w = Inches(1.45)
box_h = Inches(1.1)
for i, (col, label, sub) in enumerate(flow_steps):
    row = i // 4
    col_idx = i % 4
    x = Inches(0.35) + col_idx * (box_w + Inches(0.12))
    y = Inches(1.4) + row * (box_h + Inches(0.55))
    # Arrow
    if i < 7:
        ax = x + box_w
        ay = y + box_h / 2
        add_rect(slide, ax, ay - Pt(2), Inches(0.12), Pt(4), fill_rgb=C_MID_GRAY)
    # Box
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run1 = p.add_run(); run1.text = label; run1.font.size = Pt(10)
    run1.font.bold = True; run1.font.color.rgb = C_WHITE; run1.font.name = "Calibri"
    run2 = p.add_run(); run2.text = f"\n{sub}"
    run2.font.size = Pt(8); run2.font.color.rgb = C_WHITE; run2.font.name = "Calibri"
    run2.font.italic = True

# Tool ecosystem
tools = [
    ("get_auth_context()", "Fetches user BUKRS/WERKS/EKORG scope from SAP Auth Profile"),
    ("search_schema()", "Qdrant retrieval of DDIC metadata — tables, columns, descriptions"),
    ("get_sql_patterns()", "Few-shot proven SQL templates — MANDT, LOEVM, join idioms"),
    ("traverse_graph()", "NetworkX shortest-path between SAP tables across modules"),
    ("execute_hana_sql()", "Runs validated SQL on HANA, returns results or error for retry"),
]
y_start = Inches(4.15)
add_textbox(slide, "Agent Tool Ecosystem",
            Inches(0.35), y_start, Inches(12.6), Inches(0.3),
            font_size=11, bold=True, color=C_TEXT_DARK)
for i, (tool, desc) in enumerate(tools):
    y = y_start + Inches(0.35) + i * Inches(0.5)
    add_rect(slide, Inches(0.35), y + Inches(0.06), Pt(5), Pt(5),
             fill_rgb=PILLAR_COLORS[i // 2 + (1 if i > 2 else 0)] if i < 5 else C_ACCENT_BLUE)
    add_textbox(slide, tool,
                Inches(0.5), y, Inches(2.5), Inches(0.42),
                font_size=10, bold=True, color=C_ACCENT_BLUE)
    add_textbox(slide, desc,
                Inches(3.0), y, Inches(9.9), Inches(0.42),
                font_size=10, color=C_TEXT_DARK)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – SECTION DIVIDER: PILLAR 3
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_divider(slide, 3, "Schema RAG", "Pillar 3 — Dynamic SAP DDIC Metadata Retrieval")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – SCHEMA RAG DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Schema RAG — DDIC to Vector Store Pipeline",
           "From 90,000 SAP tables to lean, retrieved context windows")

# Pipeline boxes
pipeline = [
    (C_ACCENT_BLUE, "SAP DDIC\nExtraction",
     "DD02L / DD03L / DD08L\nPyRFC / OData"),
    (RGBColor(0x00,0x5F,0xA3), "Document\nEnrichment",
     "JSON schema docs:\ntables + columns\n+ join paths"),
    (RGBColor(0x00,0x48,0x88), "Embedding\n& Index",
     "text-embedding-3-large\n→ Qdrant collection\nsap_schemas"),
    (C_ORANGE, "Runtime\nRetrieval",
     "search_schema(query,\nmodule, top_k=3)"),
]

for i, (col, label, sub) in enumerate(pipeline):
    x = Inches(0.4) + i * Inches(3.18)
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(1.5), Inches(2.9), Inches(1.4))
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run1 = tf.paragraphs[0].add_run()
    run1.text = label; run1.font.size = Pt(12); run1.font.bold = True
    run1.font.color.rgb = C_WHITE; run1.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run(); run2.text = f"\n{sub}"
    run2.font.size = Pt(9.5); run2.font.color.rgb = C_WHITE; run2.font.name = "Calibri"
    if i < 3:
        add_rect(slide, x + Inches(2.9), Inches(2.1), Inches(0.28), Pt(4),
                 fill_rgb=C_MID_GRAY)

# Domain coverage table
domains = [
    ["Domain", "Core Tables", "Vector Search Filter"],
    ["Material (MM)", "MARA, MARC, MARD, MBEW, MVKE", "module=MM"],
    ["Business Partner", "BUT000, LFA1, KNA1, ADRC, BUT100", "module=BP"],
    ["Purchasing", "EINA, EINE, EORD, EKKO, EKPO", "module=MM-PUR"],
    ["Sales & Distribution", "VBAK, VBAP, KONV, KNVV, TVRO", "module=SD"],
    ["Finance / CO", "SKA1, SKB1, CSKS, CEPC, ANLA", "module=FI"],
    ["Quality Mgmt", "QMAT, QPAC, PLMK, MAPL", "module=QM"],
    ["Project System", "PROJ, PRPS, AFVC", "module=PS"],
]
tbl2 = slide.shapes.add_table(
    len(domains), 3,
    Inches(0.35), Inches(3.1), Inches(12.6), Inches(3.6)
).table
for ci, cw in enumerate([Inches(3.2), Inches(6.0), Inches(3.4)]):
    tbl2.columns[ci].width = cw
for ri, row_data in enumerate(domains):
    for ci, val in enumerate(row_data):
        tbl2.cell(ri, ci).text = val
table_style(tbl2, C_ORANGE, C_LIGHT_GRAY)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – SECTION DIVIDER: PILLAR 4
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_divider(slide, 4, "SQL RAG", "Pillar 4 — Proven SAP SQL Pattern Injection")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – SQL RAG DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "SQL RAG — Few-Shot Pattern Library",
           "Grounding LLMs in proven SAP SQL idioms instead of generating from scratch")

# Left: why it matters
reasons = [
    "MANDT (Client) filter — LLM hallucination target #1",
    "Module-specific deletion flags — LOEVM, LVORM, LVOMA",
    "Header/Item join patterns — VBAK→VBAP→KONV",
    "Index-aware join order — BUKRS first on BSEG/MSEG",
    "Business logic — ADD_DAYS(ZFBDT, ZBD3T) for overdue",
]
add_rect(slide, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.8),
         fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0,0xDA,0xE4), line_width_pt=0.5)
add_rect(slide, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.45),
         fill_rgb=C_RED)
add_textbox(slide, "Critical SAP Idioms LLMs Miss",
            Inches(0.5), Inches(1.38), Inches(5.6), Inches(0.38),
            font_size=11, bold=True, color=C_WHITE)
for i, r in enumerate(reasons):
    y = Inches(1.95) + i * Inches(0.9)
    add_rect(slide, Inches(0.5), y + Inches(0.05), Pt(5), Pt(5), fill_rgb=C_RED)
    add_textbox(slide, r,
                Inches(0.7), y, Inches(5.3), Inches(0.8),
                font_size=10.5, color=C_TEXT_DARK)

# Right: SQL library seed entries
sql_lib = [
    ["Query ID", "Intent / Business Question", "Tables", "Domain"],
    ["MAT-001", "Show material base details and weight", "MARA, MAKT", "MM"],
    ["MAT-003", "Current stock for material in plant", "MARD, MARC", "MM"],
    ["SLS-001", "Sales order header and status", "VBAK, VBUK", "SD"],
    ["SLS-002", "Items and pricing for sales order", "VBAP, KONV", "SD"],
    ["FIN-001", "G/L account details for company", "SKA1, SKB1", "FI"],
    ["FIN-002", "Open vendor items (overdue)", "BSIK, LFA1", "FI"],
    ["PUR-001", "Open POs for a specific vendor", "EKKO, EKPO, LFA1", "MM"],
    ["BP-001", "Vendor bank and payment terms", "LFA1, LFB1, LFBK", "BP"],
]
tbl3 = slide.shapes.add_table(
    len(sql_lib), 4,
    Inches(6.45), Inches(1.35), Inches(6.55), Inches(5.8)
).table
for ci, cw in enumerate([Inches(0.9), Inches(2.8), Inches(1.7), Inches(1.15)]):
    tbl3.columns[ci].width = cw
for ri, row_data in enumerate(sql_lib):
    for ci, val in enumerate(row_data):
        tbl3.cell(ri, ci).text = val
table_style(tbl3, C_RED, C_LIGHT_GRAY)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – SECTION DIVIDER: PILLAR 5
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
section_divider(slide, 5, "Graph RAG", "Pillar 5 — Cross-Module Relationship Mapping via NetworkX")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – GRAPH RAG DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Graph RAG — NetworkX Multi-Hop Traversal",
           "Finding join paths that Schema RAG cannot reach")

# Left: graph concept
add_rect(slide, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.8),
         fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0,0xDA,0xE4), line_width_pt=0.5)
add_rect(slide, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.45),
         fill_rgb=C_TEAL)
add_textbox(slide, "NetworkX Graph Construction",
            Inches(0.5), Inches(1.38), Inches(5.6), Inches(0.38),
            font_size=11, bold=True, color=C_WHITE)

# Node boxes
nodes = [
    (C_TEAL,    "LFA1",    "Vendor Master",     Inches(0.6),  Inches(4.3)),
    (RGBColor(0x00,0x78,0x6C), "EINA",  "Purchasing\nInfo Record", Inches(2.2), Inches(3.5)),
    (RGBColor(0x00,0x60,0x55), "MARA",  "Material Master", Inches(3.8), Inches(4.3)),
    (RGBColor(0x00,0x48,0x40), "CSKS",  "Cost Center\nMaster", Inches(5.4), Inches(3.0)),
]
for col, name, desc, nx, ny in nodes:
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                nx, ny, Inches(1.4), Inches(0.95))
    bx.fill.solid(); bx.fill.fore_color.rgb = col; bx.line.fill.background()
    tf = bx.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1 = tf.paragraphs[0].add_run()
    r1.text = name; r1.font.size = Pt(12); r1.font.bold = True
    r1.font.color.rgb = C_WHITE; r1.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(8)
    r2.font.color.rgb = C_WHITE; r2.font.name = "Calibri"

# Edges (lines with labels)
edge_data = [
    (Inches(1.9), Inches(4.7), Inches(0.4), Pt(3), "LIFNR", Inches(1.4), Inches(4.95)),
    (Inches(3.5), Inches(4.7), Inches(0.4), Pt(3), "MATNR", Inches(4.15), Inches(4.95)),
    (Inches(2.95), Inches(3.55), Inches(0.3), Pt(3), "MATNR", Inches(3.55), Inches(3.85)),
    (Inches(4.65), Inches(3.95), Inches(0.3), Pt(3), "KOSGR", Inches(5.05), Inches(3.45)),
]
for ex, ey, ew, eh, lbl, lx, ly in edge_data:
    add_rect(slide, ex, ey, ew, eh, fill_rgb=C_TEAL)
    add_textbox(slide, lbl, lx, ly, Inches(0.6), Inches(0.25),
                font_size=8, bold=True, italic=True, color=C_TEAL,
                align=PP_ALIGN.CENTER)

# Right: cross-module paths
paths = [
    ("Procurement",  "Material → Vendor",      "MARA → EINA → LFA1"),
    ("Sales",        "Material → Customer",     "MARA → KNMT → KNA1"),
    ("Costing",      "Material → Cost Center",  "MARA → MARC → MBEW → CKMLHD → CSKS"),
    ("Quality",      "Material → Inspection",  "MARA → MAPL → PLMK"),
    ("Project",      "Material → WBS",          "MARA → RESB → PRPS"),
    ("Finance",      "Vendor → G/L Account",   "LFA1 → BSIK → HKONT → SKA1"),
]
add_rect(slide, Inches(6.45), Inches(1.35), Inches(6.55), Inches(5.8),
         fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0,0xDA,0xE4), line_width_pt=0.5)
add_rect(slide, Inches(6.45), Inches(1.35), Inches(6.55), Inches(0.45),
         fill_rgb=C_TEAL)
add_textbox(slide, "Cross-Module Join Paths",
            Inches(6.6), Inches(1.38), Inches(6.2), Inches(0.38),
            font_size=11, bold=True, color=C_WHITE)

for i, (integration, query, path) in enumerate(paths):
    y = Inches(1.95) + i * Inches(0.85)
    add_rect(slide, Inches(6.55), y + Inches(0.05), Pt(5), Pt(5),
             fill_rgb=C_TEAL)
    add_textbox(slide, integration,
                Inches(6.75), y, Inches(1.8), Inches(0.35),
                font_size=10, bold=True, color=C_TEAL)
    add_textbox(slide, query,
                Inches(6.75), y + Inches(0.3), Inches(2.5), Inches(0.35),
                font_size=9.5, italic=True, color=C_TEXT_DARK)
    add_textbox(slide, path,
                Inches(9.3), y, Inches(3.5), Inches(0.7),
                font_size=9, color=C_MID_GRAY, italic=True)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – UNIFIED ARCHITECTURE STACK
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Unified 5-Pillar Stack", "End-to-end flow from user query to secured answer")

layers_stack = [
    (C_DARK_NAVY,  "User Interface",         "Streamlit / Gradio — Natural Language Chat, Session History, Export"),
    (C_NAVY,       "Pillar 1: Role-Aware RAG","AuthContext (BUKRS/WERKS) + SQL Validator + Response Masker"),
    (RGBColor(0x28,0x55,0x8A), "Pillar 2: Agentic RAG",   "ReAct Loop — Intent Routing → Tool Selection → Self-Correction"),
    (RGBColor(0x00,0x7A,0xC2), "Pillars 3+4+5",           "Schema RAG + SQL RAG + Graph RAG — Specialised Retrieval Mesh"),
    (RGBColor(0xF5,0xA6,0x23), "SQL Generation & Validation","LLM → HANA SQL → sqlglot AST check → AuthConstraint gate"),
    (RGBColor(0x00,0x96,0x88), "SAP S/4 HANA Database",   "Trusted execution — MANDT-enforced, index-aware SQL"),
    (C_GREEN,      "Audit & Response Layer",  "Immutable ledger — user, SQL hash, auth result, rows returned"),
]

for i, (col, layer, desc) in enumerate(layers_stack):
    h = Inches(0.72)
    y = Inches(1.3) + i * h
    indent = i * Inches(0.35)
    add_rect(slide, Inches(0.35) + indent, y, W - Inches(0.7) - indent, h - Pt(3),
             fill_rgb=col)
    add_textbox(slide, layer,
                Inches(0.5) + indent, y + Pt(4), Inches(3.5), h - Pt(8),
                font_size=11, bold=True, color=C_WHITE)
    add_textbox(slide, desc,
                Inches(4.0) + indent, y + Pt(4), W - Inches(4.5) - indent, h - Pt(8),
                font_size=10, color=C_WHITE, italic=True)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – MASTER DATA DOMAINS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Enterprise Master Data Coverage", "18 SAP modules supported at launch")

doms = [
    ("Business Partner", "Cross-Module", "Vendors, Customers, Roles, Addresses", "BUT000, LFA1, KNA1, ADRC"),
    ("Material Master",  "MM / SD / PP", "Materials, Products, Services, Plant Data", "MARA, MARC, MARD, MBEW"),
    ("Purchasing",       "MM-PUR",       "Purchasing Info Records, Source Lists, Quotas", "EINA, EINE, EORD, EQUK"),
    ("Sales & Dist.",    "SD",           "Pricing, Customer-Material, Routes", "KONP, KNMT, TVRO"),
    ("Warehouse Mgmt",   "EWM / WM",     "Storage Bins, Types, Handling Units", "LAGP, LQUA, VEKP"),
    ("Quality Mgmt",     "QM",           "Inspection Plans, Master Characteristics", "MAPL, PLMK, QINF"),
    ("Project System",  "PS",           "WBS Elements, Networks, Project Defs", "PRPS, PROJ, AFVC"),
    ("Transportation",   "TM / LE-TRA",  "Freight Forwarders, Transportation Zones", "/SCMTMS/D_TORROT, VTTK"),
    ("Customer Service", "CS",           "Service Masters, Warranties, Contracts", "ASMD, BGMK, VBAK"),
    ("EHS",             "EHS",          "Substance Data, Dangerous Goods, Specs", "ESTRH, ESTVH, DGTMD"),
    ("Variant Config",  "LO-VC",        "Characteristics, Classes, Config Profiles", "CABN, KLAH, CUOBJ"),
    ("Real Estate",     "RE-FX",        "Architectural Views, RE Contracts", "VICNCN, VIBDAO"),
    ("GTS",             "GTS",          "Customs, Commodity Codes, Sanctioned Parties", "/SAPSLL/PNTPR"),
    ("IS-OIL",          "IS-OIL",       "Silo/Tank Data, Joint Venture Accounting", "OIB_A04, T8JV"),
    ("IS-Retail",       "IS-R",         "Article Master, Site Master, Assortments", "MARA, T001W, WRS1"),
    ("IS-Utilities",    "IS-U",         "Device Locations, Installations, Connection Objects", "EGERR, EANL, EVBS"),
    ("IS-Health",       "IS-H",         "Patients, Business Partners, Cases", "NPAT, NBEW, NPNZ"),
    ("Taxation India",   "CIN / GST",    "HSN/SAC Codes, GSTIN, Vendor GST Details", "J_1IG_HSN_SAC, J_1BBRANCH"),
]

tbl4 = slide.shapes.add_table(
    len(doms) + 1, 4,
    Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.8)
).table
for ci, cw in enumerate([Inches(2.8), Inches(1.8), Inches(4.2), Inches(3.8)]):
    tbl4.columns[ci].width = cw

hdr = ["Domain", "SAP Module", "Key Objects", "Core Tables"]
for ci, h in enumerate(hdr):
    tbl4.cell(0, ci).text = h
table_style(tbl4, C_NAVY, C_LIGHT_GRAY)

for ri, (dom, mod, objs, tabs) in enumerate(doms):
    row = tbl4.rows[ri + 1]
    for ci, val in enumerate([dom, mod, objs, tabs]):
        row.cells[ci].text = val

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – IMPLEMENTATION PHASES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header_bar(slide, "Implementation Roadmap", "Build foundations fast, add cross-module intelligence incrementally")

phases = [
    (C_ACCENT_BLUE, "Phase 1", "Weeks 1–4",
     "Isolated Master Data Foundations",
     ["Schema RAG + SQL RAG for BP, Material, Finance",
      "Role-Aware RAG core — AuthContext service",
      "Agentic Orchestrator — basic ReAct loop",
      "Fast win: Chat with single-domain tables"]),
    (RGBColor(0x00,0x7A,0xC2), "Phase 2", "Weeks 5–8",
     "Cross-Module via Graph RAG",
     ["NetworkX graph construction from DD08L",
      "Graph RAG for MM→BP→FI join paths",
      "Multi-hop Q&A: MARA→EINA→LFA1→CSKS",
      "CDS View abstraction layer"]),
    (RGBColor(0x00,0x60,0x99), "Phase 3", "Weeks 9–11",
     "Transactional Data Integration",
     ["Extend to PO, Invoice, Material Docs, JEs",
      "High-volume table support: EKKO/EKPO, BSEG",
      "SQL Library expansion — aggregation patterns",
      "Response masking refinement"]),
    (C_DARK_NAVY, "Phase 4", "Weeks 12–14",
     "Production Polish & Audit",
     ["Immutable audit ledger — SOX ready",
      "UAT with distinct persona roles",
      "Performance benchmarking & optimization",
      "Production deployment & monitoring"]),
]

for i, (col, phase, timing, title, bullets) in enumerate(phases):
    x = Inches(0.3) + i * Inches(3.25)
    add_rect(slide, x, Inches(1.35), Inches(3.1), Inches(5.8),
             fill_rgb=C_WHITE, line_rgb=RGBColor(0xD0,0xDA,0xE4), line_width_pt=0.5)
    add_rect(slide, x, Inches(1.35), Inches(3.1), Inches(1.05), fill_rgb=col)
    add_textbox(slide, phase,
                x + Inches(0.1), Inches(1.38), Inches(1.4), Inches(0.3),
                font_size=9, bold=True, color=C_WHITE)
    add_textbox(slide, timing,
                x + Inches(1.5), Inches(1.38), Inches(1.4), Inches(0.3),
                font_size=9, color=C_ACCENT_GOLD)
    add_textbox(slide, title,
                x + Inches(0.1), Inches(1.68), Inches(2.9), Inches(0.6),
                font_size=11, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        y = Inches(2.55) + j * Inches(1.05)
        add_rect(slide, x + Inches(0.12), y + Inches(0.06), Pt(5), Pt(5),
                 fill_rgb=col)
        add_textbox(slide, b,
                    x + Inches(0.25), y, Inches(2.75), Inches(0.95),
                    font_size=9.5, color=C_TEXT_DARK)

footer_bar(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 – CLOSING / CTA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)

# Right accent panel
add_rect(slide, Inches(9.2), 0, Inches(4.13), H, fill_rgb=C_NAVY)
add_rect(slide, Inches(9.13), 0, Pt(6), H, fill_rgb=C_ACCENT_BLUE)

# Gold line
add_rect(slide, Inches(0.55), Inches(1.6), Inches(3.5), Pt(3),
         fill_rgb=C_ACCENT_GOLD)

add_textbox(slide, "From Vendor Chatbot\nto Enterprise Intelligence",
            Inches(0.5), Inches(1.75), Inches(8.4), Inches(1.5),
            font_size=34, bold=True, color=C_WHITE)

add_textbox(slide,
    "The 5-Pillar RAG Architecture transforms how enterprises interact with SAP S/4 HANA — "
    "securely, accurately, and across every master data domain.",
            Inches(0.5), Inches(3.4), Inches(8.4), Inches(1.0),
            font_size=14, color=C_LIGHT_BLUE, italic=True)

# Key stats
stats = [
    ("18+", "SAP Modules\nCovered"),
    ("5", "RAG Pillars\nOrchestrated"),
    ("90K+", "Tables\nRetrieved"),
    ("3-Layer", "Security\nEnforcement"),
]
for i, (num, lbl) in enumerate(stats):
    x = Inches(0.5) + i * Inches(2.1)
    add_textbox(slide, num,
                x, Inches(4.6), Inches(1.8), Inches(0.9),
                font_size=36, bold=True, color=C_ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, lbl,
                x, Inches(5.45), Inches(1.8), Inches(0.6),
                font_size=10, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

# Right panel CTA
add_textbox(slide, "Next Steps",
            Inches(9.45), Inches(1.5), Inches(3.5), Inches(0.5),
            font_size=18, bold=True, color=C_WHITE)
next_steps = [
    "Deploy Schema RAG for your core domain (BP or MM)",
    "Connect SAP HANA via PyRFC / OData",
    "Seed SQL Library with your first 10 validated queries",
    "Onboard 3 pilot roles for AuthContext testing",
    "Schedule a demo with your SAP Basis team",
]
for i, s in enumerate(next_steps):
    y = Inches(2.15) + i * Inches(0.9)
    add_rect(slide, Inches(9.45), y + Inches(0.08), Pt(5), Pt(5),
             fill_rgb=C_ACCENT_BLUE)
    add_textbox(slide, s,
                Inches(9.65), y, Inches(3.3), Inches(0.8),
                font_size=10.5, color=C_WHITE)

footer_bar(slide, "Know Your SAP Masters — Confidential")


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = "C:/Users/vishnu/.openclaw/workspace/SAP_HANA_LLM_VendorChatbot/Know_Your_SAP_Masters_5Pillar_RAG.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
