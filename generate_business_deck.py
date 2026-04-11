"""
Know Your SAP Masters — Deck 2: Business Team (Procurement / Finance / AP)
SAP Joule vs. Know Your SAP Masters — Value Proposition
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Palette ──────────────────────────────────────────────
C_DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x3E)
C_NAVY        = RGBColor(0x1A, 0x37, 0x6E)
C_ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xC2)
C_ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF5)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF7)
C_MID_GRAY    = RGBColor(0x8A, 0x9B, 0xA8)
C_TEXT_DARK   = RGBColor(0x1C, 0x2D, 0x4A)
C_GREEN       = RGBColor(0x00, 0xB0, 0x6A)
C_TEAL        = RGBColor(0x00, 0x96, 0x88)
C_ORANGE      = RGBColor(0xE8, 0x7D, 0x1A)
C_RED         = RGBColor(0xCC, 0x29, 0x2A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, lw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = Pt(lw) if lw else 0
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, sz=11, bold=False, italic=False,
        color=C_TEXT_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.1), fill_rgb=C_DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), W, Pt(4), fill_rgb=C_ACCENT_GOLD)
    txt(slide, title, 0.4, 0.12, W - 0.8, 0.65, sz=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.72, W - 0.8, 0.35, sz=13, color=C_ACCENT_GOLD)


def footer(slide, text="Know Your SAP Masters — Business Team Deck  |  Confidential"):
    add_rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill_rgb=C_DARK_NAVY)
    txt(slide, text, 0.3, H - Inches(0.28), W - 0.6, 0.26, sz=8, color=C_MID_GRAY)


def section(slide, num, title, subtitle=""):
    add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
    txt(slide, f"0{num}", 0.5, 1.2, 4, 3.5, sz=160, bold=True,
        color=RGBColor(0x1A, 0x37, 0x6E))
    add_rect(slide, 4.5, 2.9, 0.06, 1.8, fill_rgb=C_ACCENT_GOLD)
    txt(slide, title, 4.8, 2.7, 8, 1.2, sz=38, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 4.8, 4.0, 8, 0.6, sz=16, color=C_ACCENT_GOLD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
add_rect(slide, 9.2, 0, 4.13, H, fill_rgb=C_NAVY)
add_rect(slide, 9.13, 0, Pt(6), H, fill_rgb=C_ACCENT_GOLD)

lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 1.5, 3.2, 0.42)
lbl.fill.solid(); lbl.fill.fore_color.rgb = C_ACCENT_GOLD; lbl.line.fill.background()
tf = lbl.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "BUSINESS TEAM — PROCUREMENT & FINANCE"
r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C_WHITE

txt(slide, "Know Your\nSAP Masters", 0.5, 2.1, 8.4, 2.2, sz=52, bold=True, color=C_WHITE)
txt(slide, "Vendor Master Intelligence — Built for Procurement & Finance Teams",
    0.5, 4.35, 8.4, 0.55, sz=18, color=C_ACCENT_GOLD)
add_rect(slide, 0.5, 5.05, 3.5, Pt(3), fill_rgb=C_ACCENT_BLUE)
txt(slide, "Ask in plain English. Get SAP SQL intelligence. Stay in control.",
    0.5, 5.2, 8.4, 0.45, sz=13, italic=True, color=C_MID_GRAY)

features = [
    ("Query Vendor Data", "LFA1, LFB1, LFBK, ADRC"),
    ("Negotiation Insights", "CLV, PSI, Churn Risk, BATNA"),
    ("20-Year History", "Temporal Engine + Economic Cycles"),
    ("Payment Intelligence", "Overdue analysis, BSIK/BSAK"),
    ("QM Long-Text", "20yr inspection complaints"),
]
for i, (feat, tables) in enumerate(features):
    y = 1.0 + i * 1.1
    add_rect(slide, 9.5, y, 3.4, 0.9, fill_rgb=C_NAVY)
    txt(slide, feat, 9.65, y + 0.05, 3.1, 0.35, sz=10, bold=True, color=C_ACCENT_GOLD)
    txt(slide, tables, 9.65, y + 0.42, 3.1, 0.35, sz=9, color=C_MID_GRAY, italic=True)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM FOR BUSINESS TEAMS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "The Problem: Business Teams Can't Self-Serve SAP Data",
       "Current reality: procurement and finance depend on Basis for every query")

pain_points = [
    ("Slow Report Cycles",
     "Weekly or monthly procurement reports. By the time data arrives, it's stale. Every 'quick question' to Basis takes 3-5 business days."),
    ("No Vendor Intelligence",
     "SAP has 20 years of vendor history — delivery failures, quality issues, price trends, payment behavior. Nobody can synthesize it fast enough to use in a negotiation."),
    ("Incomplete Information at Negotiation",
     "Buyers walk into contract negotiations without CLV, churn risk, PSI score, or competitor pricing history. They lose negotiating leverage."),
    ("No Plain-English Access",
     "SAP transactions require transaction codes (ME1L, FK03, XK03) and field-level knowledge that business users shouldn't need to memorize."),
    ("QM Complaints Buried in Text",
     "Quality complaint long-text (QMEL-LONGTEXT) has 20 years of inspection notes, defect descriptions, and root cause analysis — inaccessible without a Basis report."),
]

for i, (title, desc) in enumerate(pain_points):
    y = 1.35 + i * 1.18
    add_rect(slide, 0.4, y, 0.08, 0.9, fill_rgb=C_RED)
    txt(slide, title, 0.65, y, 4.0, 0.35, sz=12, bold=True, color=C_RED)
    txt(slide, desc, 0.65, y + 0.33, 12.0, 0.58, sz=10.5, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION: WHAT BUSINESS TEAMS GET
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "The Solution: Vendor Master Intelligence in Plain English",
       "Know Your SAP Masters gives business teams direct, secure, intelligent access to SAP")

cards = [
    (C_ACCENT_BLUE, "1", "Plain English Queries",
     ["Ask: 'Show me vendors in Germany with overdue invoices'",
      "Ask: 'What materials did Vendor X supply in Q3 2024?'",
      "Ask: 'Which vendors have the best delivery reliability?'",
      "No transaction codes. No Basis ticket required."]),
    (C_GREEN, "2", "Vendor Negotiation Intelligence",
     ["20-year vendor CLV, PSI score, churn risk",
      "Competitor pricing history from EINA/EINE",
      "Delivery reliability: on-time % over 3 years",
      "Payment behavior: early payer or chronic late?"]),
    (C_TEAL, "3", "Procurement Performance Analytics",
     ["Spend by vendor, category, plant, FY",
      "Supplier Performance Index (composite 0-100)",
      "Open PO risk: vendor financial health signals",
      "Economic cycle tagging: COVID, inflation impact"]),
    (C_ORANGE, "4", "Finance: AP & Payment Intelligence",
     ["Overdue vendor items with aging buckets",
      "Payment terms analysis (LFB1-ZTERM)",
      "Bank account changes (audit risk detection)",
      "Invoice dispute history and resolution time"]),
]

for i, (col, num, title, bullets) in enumerate(cards):
    row = i // 2; ci = i % 2
    x = 0.35 + ci * 6.5; y = 1.35 + row * 2.95
    add_rect(slide, x, y, 6.3, 2.8, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, x, y, 6.3, 0.55, fill_rgb=col)
    add_rect(slide, x + 0.1, y + 0.08, 0.4, 0.4, fill_rgb=RGBColor(0xFF, 0xFF, 0xFF))
    txt(slide, num, x + 0.1, y + 0.08, 0.4, 0.4, sz=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, title, x + 0.6, y + 0.1, 5.5, 0.4, sz=13, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = y + 0.65 + j * 0.52
        add_rect(slide, x + 0.2, yy + 0.06, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.4, yy, 5.7, 0.45, sz=10, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — REAL BUSINESS QUESTIONS ANSWERED
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Real Questions Business Teams Actually Ask",
       "Know Your SAP Masters translates natural language into SAP SQL — automatically")

questions = [
    ("Procurement",   "Which vendors in Germany have the highest overdue invoice balance?",     "EKKO + BSIK + LFA1 + ADRC",  C_ACCENT_BLUE),
    ("Procurement",   "Show me the Supplier Performance Index for our top 10 vendors",          "EKKO + EKET + MSEG + QALS",   C_GREEN),
    ("Finance",       "What is the payment history trend for Vendor X over the last 3 years?",  "BSIK + BSEG + LFA1 (temporal)", C_TEAL),
    ("Negotiation",   "What price increases has Vendor X achieved in the last 5 years?",        "EINA + KONP + AQP1 (temporal)", C_ORANGE),
    ("Quality",       "Summarize all QM complaints for Material Y in the last 24 months",       "QALS + QMEL (semantic)",       C_ACCENT_BLUE),
    ("Finance",       "Which vendors changed their bank account details in the last 90 days?",  "LFBK + CDHDR + CDPOS",         C_RED),
    ("Procurement",   "What is the current stock level for Material X across all plants?",     "MARD + MARC + MBEW",           C_GREEN),
    ("Negotiation",   "Show me the CLV, churn risk, and BATNA for Customer Y",                 "VBAP + BSID + KNVV + AQFE",   C_TEAL),
]

for i, (dept, question, tables, col) in enumerate(questions):
    row = i // 2; ci = i % 2
    x = 0.35 + ci * 6.5; y = 1.35 + row * 1.45
    add_rect(slide, x, y, 6.3, 1.35, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, x, y, 1.1, 1.35, fill_rgb=col)
    txt(slide, dept, x, y, 1.1, 1.35, sz=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, question, x + 1.2, y + 0.08, 4.9, 0.7, sz=10.5, color=C_TEXT_DARK)
    txt(slide, tables, x + 1.2, y + 0.82, 4.9, 0.4, sz=9, italic=True, color=C_MID_GRAY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — NEGOTIATION BRIEFING (PHASE 8)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Negotiation Briefing Generator — Your AI Strategy Assistant",
       "Before every vendor negotiation: full intelligence brief in under 30 seconds")

brief = [
    (C_ACCENT_GOLD, "CLV Tier",
     ["GOLD: >₹50L annual spend | SILVER: ₹10-50L | BRONZE: <₹10L",
      "Based on 3-year revenue from VBAP-NETWR",
      "Used to set negotiating leverage baseline"]),
    (C_GREEN, "PSI Score (Supplier Performance Index)",
     ["Composite 0-100: Delivery (40%) + Quality (35%) + Price (25%)",
      "Delivery: EKKO/EKET/MSEG on-time % over 12 months",
      "Quality: QALS UD code analysis (0=accepted, 1=deviation)"]),
    (C_TEAL, "Churn Risk Signal",
     ["TRIGGER: >90-day gap in purchase orders → flagged as CHURN RISK",
      "TRIGGER: >15% YoY price increase accepted → vendor has pricing power",
      "TRIGGER: Quality reject rate >5% → vendor dependency risk"]),
    (C_ORANGE, "BATNA & Tactics",
     ["Top tactic library: 8 strategies based on PSI + churn + CLV",
      "If PSI > 75 AND churn = LOW → 'You have pricing power — hold firm'",
      "If PSI < 40 AND churn = HIGH → 'Vendor needs you more than you need them'"]),
]

for i, (col, name, bullets) in enumerate(brief):
    x = 0.35 + (i % 2) * 6.5; y = 1.35 + (i // 2) * 2.95
    add_rect(slide, x, y, 6.3, 2.8, fill_rgb=C_WHITE, line_rgb=col, lw=1.5)
    add_rect(slide, x, y, 6.3, 0.5, fill_rgb=col)
    txt(slide, name, x + 0.15, y + 0.08, 6.0, 0.4, sz=12, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = y + 0.6 + j * 0.55
        add_rect(slide, x + 0.2, yy + 0.06, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.4, yy, 5.7, 0.5, sz=10, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — TEMPORAL / 20-YEAR HISTORY
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "20-Year Longitudinal Corporate Memory",
       "SAP from 2005 → your competitive advantage in 2026")

use_cases = [
    ("2008 Financial Crisis Analysis",
     "Which raw material categories saw the highest vendor default rates in 2008-2010? Who were backup vendors we pivoted to?",
     "LFA1 + EKKO + EKPO + MSEG + BSEG",
     C_RED),
    ("COVID-19 Supply Chain Shock",
     "Which vendors went silent (no POs) during Mar-Jun 2020? How long did recovery take?",
     "EKKO + EKET + MSEG (temporal filter: 2020)",
     C_ORANGE),
    ("Inflation Squeeze 2022",
     "How did our top 20 vendors behave when we pushed back on price increases? Who accepted, who churned?",
     "EINA + KONP + EKKO (temporal: FY2022-FY2023)",
     C_TEAL),
    ("Post-Retirement Vendor Intelligence",
     "A senior buyer retired — all tacit knowledge about Vendor X's 2012 payment disputes is gone from human memory. It's in SAP.",
     "BSIK + LFA1 + CDHDR + CDPOS",
     C_ACCENT_BLUE),
    ("Warranty & Quality Recalls",
     "Which materials had repeated quality defects over 10 years? Are there seasonal patterns (summer batches, winter coating failures)?",
     "QALS + QMEL + MARA + MAPL",
     C_GREEN),
]

for i, (title, desc, tables, col) in enumerate(use_cases):
    y = 1.35 + i * 1.18
    add_rect(slide, 0.35, y, 0.1, 0.9, fill_rgb=col)
    txt(slide, title, 0.55, y, 4.5, 0.35, sz=11, bold=True, color=col)
    txt(slide, desc, 0.55, y + 0.35, 8.5, 0.55, sz=10, color=C_TEXT_DARK)
    txt(slide, tables, 9.2, y + 0.25, 3.8, 0.5, sz=9, italic=True, color=C_MID_GRAY)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY & ROLE-BASED ACCESS (BUSINESS CONTEXT)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Security Built In — Role-Based Access for Business Teams",
       "Every query is scoped to what your SAP role is allowed to see")

roles = [
    ("AP Clerk",       C_ACCENT_BLUE,
     ["Access: Open invoices, payment history, vendor master basic",
      "Masked: Bank account numbers (STCD1), credit limits",
      "Scope: Own company code (BUKRS) only"]),
    ("Procurement Buyer", C_GREEN,
     ["Access: POs, vendor evaluations, pricing conditions",
      "Masked: Other buyers' negotiation notes",
      "Scope: Own purchasing organization (EKORG) + plants"]),
    ("Finance Manager", C_TEAL,
     ["Access: Full AP/AR, G/L, cost centers, profit centers",
      "Masked: HR salary data, personal employee info",
      "Scope: All BUKRS in scope for role"]),
    ("Category Manager", C_ORANGE,
     ["Access: Spend analytics, vendor performance dashboards",
      "Masked: Individual employee data, internal cost details",
      "Scope: Category (MATKL) + region (REGIO) filters"]),
]

for i, (role, col, bullets) in enumerate(roles):
    row = i // 2; ci = i % 2
    x = 0.35 + ci * 6.5; y = 1.35 + row * 2.9
    add_rect(slide, x, y, 6.3, 2.75, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, x, y, 6.3, 0.5, fill_rgb=col)
    txt(slide, role, x + 0.15, y + 0.08, 6.0, 0.4, sz=12, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = y + 0.6 + j * 0.5
        add_rect(slide, x + 0.2, yy + 0.06, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.4, yy, 5.7, 0.45, sz=10, color=C_TEXT_DARK)

txt(slide, "Security enforced at 3 layers: Prompt Injection → SQL WHERE Filter → Response Masking",
    0.35, 7.1, 12.6, 0.3, sz=9, bold=True, color=C_ACCENT_BLUE, align=PP_ALIGN.CENTER)
footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — ROI & TCO
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "ROI & Total Cost of Ownership", "Business case for Know Your SAP Masters vs. status quo")

roi_items = [
    ("Basis Ticket Elimination",  "3-5 days per request",    "5-10 requests/week",    "240-520 hrs/yr saved",    C_GREEN),
    ("Negotiation Win Rate",       "Avg price increase: 2-3%", "With intelligence: 4-6%", "+2-4% per contract",         C_ACCENT_BLUE),
    ("QM Complaint Resolution",   "Avg: 12 days to close",    "Semantic search: <1 day", "11+ days faster resolution", C_TEAL),
    ("Audit Risk Reduction",       "Bank change fraud: 2-3/yr", "Auto-detection: real-time",  "Near-zero undetected fraud", C_RED),
    ("Report Generation",          "Weekly manual: 4-6 hrs",   "Self-service: 2-3 mins",  "~200 hrs/yr recovered",    C_ORANGE),
]

# Header
add_rect(slide, 0.35, 1.35, 12.6, 0.5, fill_rgb=C_DARK_NAVY)
headers = ["Benefit", "Baseline (Status Quo)", "With KYSM", "Estimated Value"]
col_xs = [0.35, 3.35, 7.0, 10.0]
col_ws = [2.95, 3.6, 2.95, 2.65]
for h, cx, cw in zip(headers, col_xs, col_ws):
    txt(slide, h, cx + 0.05, 1.37, cw, 0.45, sz=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, (benefit, baseline, with_yysm, value, col) in enumerate(roi_items):
    y = 1.85 + ri * 1.05
    bg = C_WHITE if ri % 2 == 0 else C_LIGHT_GRAY
    add_rect(slide, 0.35, y, 12.6, 0.95, fill_rgb=bg, line_rgb=RGBColor(0xDD, 0xDD, 0xDD), lw=0.5)
    add_rect(slide, 0.35, y, 0.1, 0.95, fill_rgb=col)
    txt(slide, benefit, col_xs[0] + 0.1, y + 0.08, col_ws[0] - 0.1, 0.8, sz=10, bold=True, color=col)
    txt(slide, baseline, col_xs[1] + 0.05, y + 0.15, col_ws[1] - 0.1, 0.65, sz=10, color=C_MID_GRAY)
    txt(slide, with_yysm, col_xs[2] + 0.05, y + 0.15, col_ws[2] - 0.1, 0.65, sz=10, bold=True, color=C_GREEN)
    txt(slide, value, col_xs[3] + 0.05, y + 0.15, col_ws[3] - 0.1, 0.65, sz=10, bold=True, color=C_ACCENT_BLUE)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — SAP JOULE: COMPLEMENT NOT REPLACE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Know Your SAP Masters vs. SAP Joule — What's the Difference?",
       "Not a replacement — a domain specialist built on top of SAP")

diff = [
    ("Breadth vs. Depth",    "Joule covers all SAP modules generically",           "We focus on vendor master with 14 meta-paths and 20yr history",  C_GREEN),
    ("Query Type",            "Task automation ('run this workflow')",               "SQL generation ('show me actual SAP data')",                      C_TEAL),
    ("Negotiation Support",   "Generic CRM prompts",                                 "PSI, CLV, churn risk, BATNA, price increase analysis",           C_ACCENT_GOLD),
    ("Temporal History",      "Date-effective tables (SAP standard)",                "Temporal engine + economic cycle tagging (2008, COVID, 2022)",    C_ORANGE),
    ("Confidence",            "Joule 'reflection' — opaque",                        "6-signal breakdown — fully auditable",                           C_ACCENT_BLUE),
    ("Customization",         "SAP-managed, limited customization",                   "Open Python stack — hot-reloadable patterns & tools",            C_PURPLE := C_ACCENT_BLUE),
]

for i, (dimension, joule, ours, col) in enumerate(diff):
    y = 1.35 + i * 0.98
    bg = C_WHITE if i % 2 == 0 else C_LIGHT_GRAY
    add_rect(slide, 0.35, y, 12.6, 0.9, fill_rgb=bg, line_rgb=col, lw=0.5)
    txt(slide, dimension, 0.45, y + 0.08, 2.2, 0.75, sz=10, bold=True, color=col)
    txt(slide, joule, 2.7, y + 0.08, 4.7, 0.75, sz=9.5, color=C_MID_GRAY, italic=True)
    txt(slide, ours, 7.45, y + 0.08, 5.35, 0.75, sz=9.5, bold=True, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — IMPLEMENTATION PATH
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_LIGHT_GRAY)
header(slide, "Implementation Path — Quick Wins for Business Teams",
       "Start in 2 weeks. See results in 4 weeks. Scale in 12 weeks.")

path = [
    (C_GREEN,    "Week 1-2: Pilot",
     ["Deploy vendor master domain (LFA1 + LFB1 + LFBK + ADRC)",
      "Onboard AP Clerk and 2 procurement buyers",
      "Run first 20 queries — real business questions",
      "Baseline: how many Basis tickets does this replace?"]),
    (C_ACCENT_BLUE, "Week 3-4: Negotiation Intelligence",
     ["Activate Negotiation Briefing Generator",
      "Pull PSI scores for top 20 vendors",
      "Prepare 3 vendor negotiation briefings",
      "Measure: negotiation outcome improvement"]),
    (C_TEAL, "Week 5-8: Temporal Expansion",
     ["Activate Temporal Engine for vendor history",
      "Run 2008/COVID/inflation analysis for top vendors",
      "Feed insights into annual contract negotiations",
      "Measure: historical insight usage rate"]),
    (C_ORANGE, "Week 9-12: Full Rollout",
     ["All procurement + finance roles onboarded",
      "18 SAP master data domains active",
      "Self-service dashboard for category managers",
      "Measure: Basis ticket reduction"]),
]

for i, (col, phase, bullets) in enumerate(path):
    x = 0.35 + i * 3.2
    add_rect(slide, x, 1.35, 3.05, 5.8, fill_rgb=C_WHITE, line_rgb=col, lw=1)
    add_rect(slide, x, 1.35, 3.05, 0.55, fill_rgb=col)
    txt(slide, phase, x + 0.15, 1.4, 2.8, 0.45, sz=11, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        yy = 2.0 + j * 1.15
        add_rect(slide, x + 0.15, yy + 0.06, Pt(5), Pt(5), fill_rgb=col)
        txt(slide, b, x + 0.3, yy, 2.6, 1.05, sz=9.5, color=C_TEXT_DARK)

footer(slide)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK_NAVY)
add_rect(slide, 9.2, 0, 4.13, H, fill_rgb=C_NAVY)
add_rect(slide, 9.13, 0, Pt(6), H, fill_rgb=C_ACCENT_GOLD)
add_rect(slide, 0.55, 1.6, 3.5, Pt(3), fill_rgb=C_ACCENT_GOLD)

txt(slide, "From SAP Data Graveyard\nto Negotiation Weapon",
    0.5, 1.75, 8.4, 1.5, sz=36, bold=True, color=C_WHITE)
txt(slide,
    "Know Your SAP Masters transforms 20 years of buried SAP data into "
    "real-time intelligence for your procurement and finance teams.",
    0.5, 3.4, 8.4, 1.0, sz=14, color=C_LIGHT_BLUE, italic=True)

stats = [
    ("20yrs", "SAP Vendor\nHistory"),
    ("<30s", "Negotiation\nBrief"),
    ("5+", "Business\nRoles"),
    ("0", "Basis\nTickets"),
]
for i, (num, lbl) in enumerate(stats):
    x = 0.5 + i * 2.1
    txt(slide, num, x, 4.6, 1.8, 0.9, sz=36, bold=True, color=C_ACCENT_GOLD, align=PP_ALIGN.CENTER)
    txt(slide, lbl, x, 5.45, 1.8, 0.6, sz=10, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Your Next Steps", 9.45, 1.5, 3.5, 0.5, sz=18, bold=True, color=C_WHITE)
steps = [
    "Identify top 5 vendor negotiation questions your team asks weekly",
    "Map which SAP roles need access (AP Clerk, Buyer, Finance Manager)",
    "Schedule 30-min demo with procurement leadership",
    "Run pilot with 2 buyers and 1 AP Clerk — measure Basis ticket reduction",
    "Contact: your-SAP-masters-team",
]
for i, s in enumerate(steps):
    y = 2.15 + i * 0.9
    add_rect(slide, 9.45, y + 0.08, Pt(5), Pt(5), fill_rgb=C_ACCENT_GOLD)
    txt(slide, s, 9.65, y, 3.3, 0.8, sz=10.5, color=C_WHITE)

footer(slide, "Know Your SAP Masters — Business Team Deck")


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
out = r"C:\Users\vishnu\.openclaw\workspace\SAP_HANA_LLM_VendorChatbot\KnowYourSAPMasters_BusinessTeam_Deck.pptx"
prs.save(out)
print(f"[OK] Business Team Deck saved: {out} | Slides: {len(prs.slides)}")
